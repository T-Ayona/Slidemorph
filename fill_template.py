"""
fill_template.py  (v5 - spec driven)
=====================================

Fills a PowerPoint template with AI-generated content while PRESERVING the
manual entrance animations set up in the source deck, and fixing title
wrapping / title-vs-content overlap.

All template-specific rules (slide layout, section slide indices, per-section
slide counts, placeholder character/bullet limits, and the 6-20 slide
composition table) live in templates/<name>/spec.json. Adding a new template
means adding templates/<name>/template.pptx + templates/<name>/spec.json --
no changes to this file. This engine only knows the generic shape of a deck:

    Start (1 slide) -> Wheel (fixed count) -> Bars (variable) ->
    Hexagon (variable) -> End (1 slide)

with the exact numbers and text-fit rules coming from the spec.

Pipeline (see generate_ai_deck):
  1. load_spec(template_name)           - read templates/<name>/spec.json
  2. build_slide_sequence(count, spec)  - which template slides to clone, in order
  3. build_deck(...)                    - clone slides into a fresh .pptx
  4. generate_content(topic, spec, ...) - Gemini Stage 1: outline (titles/tags/headings)
  5. fill_placeholders(...)             - Gemini Stage 2 (adaptive bullets) + write
  6. save_output(...)

Text-fit mechanics (unchanged from earlier versions):
  - KEEP EXISTING ANIMATIONS. The fill is "structure-aware": it preserves
    every paragraph an animation targets and only rewrites the runs inside it.
  - Bullet COUNT adapts to how many bullets would wrap onto a 2nd line
    (measured with Pillow), per the spec's max_bullets_* / words_per_bullet.
  - Bullet boxes grow in height (never width) to fit; titles are widened /
    given a one-line height so they never overlap the content box, without
    ever changing font size.
  - AI_content / AI_title are renamed per-slide (AI_content_<n>, AI_title_<n>)
    only by rename_unique() (used by the dummy-fill test path) so PowerPoint
    doesn't try to morph text boxes between dissimilar slides; the AI content
    pipeline intentionally leaves names shared so Bar/Hex slides morph.
"""

import argparse
import os
import re
import sys
import copy
import json
import time
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed

# Generated content routinely contains characters the Windows console codepage
# (cp1252) cannot encode -- Greek letters in complexity notation, en-dashes, box
# glyphs. Without this, a single such character raises UnicodeEncodeError out of
# print() and aborts the run. Degrade those characters instead of the deck.
for _stream in (sys.stdout, sys.stderr):
    try:
        _stream.reconfigure(encoding="utf-8", errors="replace")
    except Exception:  # noqa: BLE001 - older/odd streams: leave them alone
        pass

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.shapes.group import GroupShape
from pptx.parts.slide import SlidePart
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.exc import PackageNotFoundError
from PIL import ImageFont

R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

# ---------------------------------------------------------------------------
# Namespaces / engine constants (generic -- apply to any template)
# ---------------------------------------------------------------------------
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS = {"a": A, "p": P}
XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"

EMU = 914400                      # EMU per inch
BULLET = "\u2022 "                # "- "
CONTENT_PT = 18                   # AI_content effective size (defaultTextStyle=1800)
TITLE_PT_DEFAULT = 32             # AI_title (Impact 3200)
LINE_FACTOR = 1.2                 # single-line spacing multiplier
FIT_SLACK_IN = 0.06               # one-line width slack for titles
TITLE_WIDTH_SLACK_IN = 0.20       # extra width when widening a title to one line
GAP_IN = 0.08                     # min gap between title bottom and content top
MIN_TOP_IN = 0.10                 # keep boxes this far from slide edges
MAX_BULLET_ITERS = 3              # regen attempts before trim

# ---- Gemini config ---------------------------------------------------------
# Tried in order; on a quota/429 OR a 404/unavailable error the run transparently
# falls back to the next model.
GEMINI_MODELS = ["gemini-2.0-flash", "gemini-flash-latest", "gemini-2.5-flash",
                 "gemini-3.5-flash", "gemini-3.6-flash", "gemini-3.1-flash-lite",
                 "gemini-flash-lite-latest", "gemini-3.5-flash-lite"]
QUOTA_RETRY_CAP = 45.0                      # max seconds to wait on a rate-limit
MAX_QUOTA_WAITS = 2                         # wait+retry attempts on the last model
API_CALL_DELAY = 1.0                        # seconds between calls (rate limiting)

# Bundled fonts live next to this file in ./fonts/, so measurements are
# identical on Windows, macOS, and Linux (no dependency on system fonts).
# Both are OFL-licensed (see fonts/OFL.txt):
#   Carlito -- metric-compatible replacement for Calibri (by Google).
#   Anton   -- condensed heavy display face used in place of Impact.
FONTS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "fonts")
_DEFAULT_FONT = "Carlito-Regular.ttf"

# Explicit typeface -> bundled file mappings. Anything not listed here (theme
# placeholders like "+mj-lt", other Windows font names, unknown typefaces) still
# resolves via font_path() below to _DEFAULT_FONT -- Pillow is NEVER handed a
# raw typeface name; it always receives a real absolute path to a bundled TTF.
FONT_FILE = {
    "Impact": "Anton-Regular.ttf",
    "Calibri Light": "Carlito-Regular.ttf",
    "Calibri": "Carlito-Regular.ttf",
    "Arial Narrow": "Carlito-Regular.ttf",
}


def font_path(typeface):
    """Absolute path to the bundled TTF for `typeface`. Guarantees a real file
    for every possible input:
      - PowerPoint theme placeholders ('+mj-lt', '+mn-lt', '+mj-ea', ...) -> Carlito
      - Windows/Mac font names not in FONT_FILE -> Carlito
      - None / empty -> Carlito
    Nothing else in the code should build a font path -- go through this."""
    if not typeface or typeface.startswith("+"):
        return os.path.join(FONTS_DIR, _DEFAULT_FONT)
    return os.path.join(FONTS_DIR, FONT_FILE.get(typeface, _DEFAULT_FONT))


_FONT_CACHE = {}
_DEFAULT_FONT_WARNED = False


def _font(typeface, pt):
    """Load and cache a font by typeface + size. Never raises: on any error it
    prints the exact path that failed and falls back to Pillow's built-in
    default (so width measurements degrade gracefully instead of crashing)."""
    global _DEFAULT_FONT_WARNED
    key = f"{typeface}-{int(round(pt))}"
    if key in _FONT_CACHE:
        return _FONT_CACHE[key]
    path = font_path(typeface)
    try:
        _FONT_CACHE[key] = ImageFont.truetype(path, int(round(pt)))
    except Exception as exc:  # noqa: BLE001 - want a total safety net
        if not _DEFAULT_FONT_WARNED:
            print(f"[font] WARNING: could not load {typeface!r} from "
                  f"{path!r} ({type(exc).__name__}: {exc}); "
                  f"falling back to Pillow's built-in default font")
            _DEFAULT_FONT_WARNED = True
        _FONT_CACHE[key] = ImageFont.load_default()
    return _FONT_CACHE[key]


def _startup_font_check():
    """One-shot diagnostic. Prints the fonts folder, its contents, and whether
    each bundled TTF is loadable. Runs at import so Streamlit Cloud logs it
    immediately on app boot -- no need to trigger a generation to see it."""
    print(f"[font] FONTS_DIR = {FONTS_DIR}")
    print(f"[font] FONTS_DIR exists = {os.path.isdir(FONTS_DIR)}")
    if os.path.isdir(FONTS_DIR):
        entries = sorted(os.listdir(FONTS_DIR))
        print(f"[font] FONTS_DIR contents ({len(entries)} item(s)): {entries}")
    for tf, fn in list(FONT_FILE.items()) + [("<default>", _DEFAULT_FONT)]:
        p = os.path.join(FONTS_DIR, fn)
        ok = os.path.isfile(p)
        loaded = False
        if ok:
            try:
                ImageFont.truetype(p, 12)
                loaded = True
            except Exception as exc:  # noqa: BLE001
                loaded = f"LOAD FAIL ({type(exc).__name__}: {exc})"
        print(f"[font]   {tf:15} -> {p}  exists={ok}  loaded={loaded}")


_startup_font_check()


def text_width_in(text, typeface, pt):
    """Rendered width of `text` in inches (font size in pt-as-px -> px==pt)."""
    return _font(typeface, pt).getlength(text) / 72.0


def bul(text):
    return BULLET + text


# Leading list markers Gemini sometimes prepends even when asked not to. The
# template already draws the bullet dot via bul(), so a marker in the text itself
# produces a visible double bullet ("* - foo" / "* * foo"). These strippers remove
# it. A marker only counts when followed by whitespace, so "-5", "->", and ">>>"
# (arrows / REPL prompts / negative numbers) are preserved.
_BULLET_MARK = "•‣◦⁃∙–—*\\-"   # • ‣ ◦ ⁃ ∙ – — * -
_BULLET_INLINE_RE = re.compile(r"^\s*[" + _BULLET_MARK + r"][ \t]+")
_BULLET_INDENT_RE = re.compile(r"^(\s*)[" + _BULLET_MARK + r">][ \t]+(.*)$")


def strip_bullet_prefix(text):
    """Remove leading list markers from a bullet string (repeatedly, to catch
    "* - foo"). Leading indentation is not meaningful for bullets, so it is
    dropped. Use this for AI_content bullets before bul() re-adds the single dot."""
    s = (text or "").strip()
    for _ in range(3):
        new = _BULLET_INLINE_RE.sub("", s, count=1)
        if new == s:
            break
        s = new.strip()
    return s


def strip_long_bullet(line):
    """Remove a single leading list marker from a long-content line WHILE
    preserving code indentation on lines that have no marker. A bulleted line
    ("    - foo") loses both marker and its indent (it was never real code
    indentation); a plain code line ("    return x") is returned untouched."""
    m = _BULLET_INDENT_RE.match(line or "")
    return m.group(2) if m else (line or "")


def wrap_to_lines(text, max_chars_per_line, max_lines):
    """Greedy word-wrap `text` into at most `max_lines` lines of <= max_chars_per_line
    each, without splitting words. If it will not fit, the last line is
    ellipsized so a word is never chopped mid-way. Returns a list of 1..max_lines
    lines (always at least one)."""
    words = (text or "").split()
    if not words:
        return [""]
    lines, cur = [], ""
    for w in words:
        trial = w if not cur else cur + " " + w
        if len(trial) <= max_chars_per_line or not cur:
            cur = trial
        else:
            lines.append(cur)
            cur = w
            if len(lines) == max_lines:         # no room for more lines
                break
    if len(lines) < max_lines:
        lines.append(cur)
    else:                                        # ran out of lines: fold the rest in
        rest = text[len(" ".join(lines)):].strip()
        tail = (lines[-1] + " " + rest).strip() if rest else lines[-1]
        lines[-1] = truncate(tail, max_chars_per_line, ellipsis=True)
    return lines[:max_lines]


# ---------------------------------------------------------------------------
# Generic helpers
# ---------------------------------------------------------------------------
def iter_all_shapes(shapes):
    for sh in shapes:
        yield sh
        if isinstance(sh, GroupShape):
            yield from iter_all_shapes(sh.shapes)


def find_shape(slide, name):
    for sh in iter_all_shapes(slide.shapes):
        if sh.name == name:
            return sh
    return None


def _xfrm(shape):
    return shape._element.find(".//a:xfrm", NS)


def geom_emu(shape):
    x = _xfrm(shape)
    off, ext = x.find("a:off", NS), x.find("a:ext", NS)
    return (int(off.get("x")), int(off.get("y")),
            int(ext.get("cx")), int(ext.get("cy")))


def set_geom(shape, ox, oy, cx, cy):
    x = _xfrm(shape)
    off, ext = x.find("a:off", NS), x.find("a:ext", NS)
    off.set("x", str(int(round(ox))));  off.set("y", str(int(round(oy))))
    ext.set("cx", str(int(round(cx))));  ext.set("cy", str(int(round(cy))))


def get_ins(shape):
    """Text-frame insets in inches (defaults: l/r=0.1, t/b=0.05)."""
    b = shape._element.find(".//a:bodyPr", NS)

    def g(attr, default):
        v = b.get(attr) if b is not None else None
        return (int(v) / EMU) if v is not None else default
    return {"l": g("lIns", 0.1), "r": g("rIns", 0.1),
            "t": g("tIns", 0.05), "b": g("bIns", 0.05)}


def run_font(shape, default_typeface, default_pt):
    """Effective (typeface, pt) of the shape's first run (or endParaRPr)."""
    el = shape._element
    rpr = el.find(".//a:r/a:rPr", NS)
    if rpr is None:
        rpr = el.find(".//a:endParaRPr", NS)
    tf, pt = default_typeface, default_pt
    if rpr is not None:
        latin = rpr.find("a:latin", NS)
        if latin is not None and latin.get("typeface"):
            tf = latin.get("typeface")
        if rpr.get("sz"):
            pt = int(rpr.get("sz")) / 100.0
    return tf, pt


# ---------------------------------------------------------------------------
# Run / paragraph surgery (formatting preserving)
# ---------------------------------------------------------------------------
def _template_rpr(txBody):
    r = txBody.find(".//a:r", NS)
    if r is not None and r.find("a:rPr", NS) is not None:
        return copy.deepcopy(r.find("a:rPr", NS))
    epr = txBody.find(".//a:endParaRPr", NS)
    if epr is not None:
        rpr = copy.deepcopy(epr)
        rpr.tag = qn("a:rPr")
        return rpr
    return None


def _para_rpr(p, txBody):
    """rPr to reuse for a paragraph's new runs (its own run first)."""
    r = p.find("a:r", NS)
    if r is not None and r.find("a:rPr", NS) is not None:
        return copy.deepcopy(r.find("a:rPr", NS))
    epr = p.find("a:endParaRPr", NS)
    if epr is not None:
        rpr = copy.deepcopy(epr)
        rpr.tag = qn("a:rPr")
        return rpr
    return _template_rpr(txBody)


def _make_run(parent, rpr, text):
    r = parent.makeelement(qn("a:r"), {})
    if rpr is not None:
        r.append(copy.deepcopy(rpr))
    t = parent.makeelement(qn("a:t"), {})
    t.text = text
    if text != text.strip():
        t.set(XML_SPACE, "preserve")
    r.append(t)
    return r


def _rebuild_para(p, lines, rpr):
    """Replace a paragraph's runs/breaks with `lines` (joined by <a:br/>),
    preserving its <a:pPr> and <a:endParaRPr>."""
    for ch in list(p):
        if ch.tag in (qn("a:r"), qn("a:br")):
            p.remove(ch)
    endpr = p.find("a:endParaRPr", NS)
    new_nodes = []
    for i, line in enumerate(lines):
        if i > 0:
            br = p.makeelement(qn("a:br"), {})
            if rpr is not None:
                br.append(copy.deepcopy(rpr))
            new_nodes.append(br)
        new_nodes.append(_make_run(p, rpr, line))
    if endpr is not None:
        for node in new_nodes:
            endpr.addprevious(node)
    else:
        for node in new_nodes:
            p.append(node)


def fill_simple(shape, lines):
    """Rebuild a whole text frame as one paragraph per line (titles/tags/subtitle).
    Preserves run + paragraph formatting; never touches <a:bodyPr>."""
    txBody = shape.text_frame._txBody
    rpr = _template_rpr(txBody)
    ppr = txBody.find(".//a:pPr", NS)
    ppr = copy.deepcopy(ppr) if ppr is not None else None
    for p in txBody.findall("a:p", NS):
        txBody.remove(p)
    for line in lines:
        p = txBody.makeelement(qn("a:p"), {})
        if ppr is not None:
            p.append(copy.deepcopy(ppr))
        p.append(_make_run(p, rpr, line))
        txBody.append(p)


# ---------------------------------------------------------------------------
# Animation-structure inspection
# ---------------------------------------------------------------------------
def animated_para_indices(slide, spid):
    """Paragraph indices the slide's timing animates for shape `spid`."""
    idx = set()
    for spt in slide._element.findall(".//p:spTgt", NS):
        if spt.get("spid") != str(spid):
            continue
        for prg in spt.findall(".//p:pRg", NS):
            idx.update(range(int(prg.get("st")), int(prg.get("end")) + 1))
    return sorted(idx)


# ---------------------------------------------------------------------------
# Content fill + fit
# ---------------------------------------------------------------------------
def content_marL_emu():
    """Left margin for the hanging indent = width of the "- " prefix (EMU)."""
    return int(round(text_width_in(BULLET, "+mj-lt", CONTENT_PT) * EMU))


def count_display_lines(text, usable_in, hang_in):
    """Greedy word-wrap line count with a hanging indent: the first line has the
    full usable width, wrapped continuation lines are narrower by `hang_in`."""
    words = text.split(" ")
    lines, cur, avail = 1, "", usable_in
    for w in words:
        trial = w if not cur else cur + " " + w
        if not cur or text_width_in(trial, "+mj-lt", CONTENT_PT) <= avail:
            cur = trial
        else:
            lines += 1
            cur = w
            avail = usable_in - hang_in     # continuation lines are indented
    return lines


def _set_pPr_hanging(p, marL_emu):
    """Give paragraph `p` a hanging indent (first line flush, wraps indented)."""
    ppr = p.find("a:pPr", NS)
    if ppr is None:
        ppr = p.makeelement(qn("a:pPr"), {})
        p.insert(0, ppr)            # <a:pPr> must be the first child
    ppr.set("marL", str(marL_emu))
    ppr.set("indent", str(-marL_emu))


def make_bullet_para(txBody, rpr, ppr_base, marL_emu, text):
    """Build one <a:p> bullet paragraph with a hanging indent."""
    p = txBody.makeelement(qn("a:p"), {})
    ppr = copy.deepcopy(ppr_base) if ppr_base is not None else p.makeelement(qn("a:pPr"), {})
    ppr.set("marL", str(marL_emu))
    ppr.set("indent", str(-marL_emu))
    p.append(ppr)
    p.append(_make_run(p, rpr, text))
    if rpr is not None:
        epr = copy.deepcopy(rpr)
        epr.tag = qn("a:endParaRPr")
        p.append(epr)
    return p


def update_animation_range(slide, spid, start, count):
    """Widen every <p:pRg> that targets shape `spid` to cover `count` paragraphs
    starting at `start` -- keeps the SAME fade effect/timing/delay, just over the
    now-split bullet paragraphs. Returns 'old -> new' for reporting."""
    changed = None
    for spt in slide._element.findall(".//p:spTgt", NS):
        if spt.get("spid") != str(spid):
            continue
        for prg in spt.findall(".//p:pRg", NS):
            old = f"{prg.get('st')}-{prg.get('end')}"
            prg.set("st", str(start))
            prg.set("end", str(start + count - 1))
            changed = f"{old} -> {start}-{start + count - 1}"
    return changed or "unchanged"


def fill_content_bullets(slide, bullets, slide_h_in, top_limit_in, marL_emu):
    """Fill AI_content with `bullets` (one hanging-indent paragraph each), fitting
    the box height and preserving any animation. Grouped fades widen to cover all
    bullet paragraphs; staggered per-bullet fades are untouched.

    Returns (final_bullets, fit_note, anim_note).
    """
    shape = find_shape(slide, "AI_content")
    spid = shape.shape_id
    animated = animated_para_indices(slide, spid)
    txBody = shape.text_frame._txBody

    # Final safety net for double bullets: the template's bul() adds the one dot,
    # so any marker still on the text here would render as "* - ...".
    bullets = [strip_bullet_prefix(b) for b in bullets]
    bullets = [b for b in bullets if b]
    staggered = len(animated) > 1
    if staggered:                   # one bullet per animated paragraph
        bullets = bullets[:len(animated)]
    hang_in = marL_emu / EMU

    # --- fit check (grow height; drop bullets only for grouped fades) ---------
    ox, oy, cx, cy = geom_emu(shape)
    ins = get_ins(shape)
    usable = cx / EMU - ins["l"] - ins["r"]
    lh = LINE_FACTOR * CONTENT_PT / 72.0

    def needed_in(bl):
        return sum(count_display_lines(bul(b), usable, hang_in) for b in bl) * lh \
            + ins["t"] + ins["b"]

    fit_note = "fits"
    if needed_in(bullets) > cy / EMU + 1e-6:
        center = oy / EMU + cy / EMU / 2.0
        top_lim = max(MIN_TOP_IN, top_limit_in)
        bot_lim = slide_h_in - MIN_TOP_IN
        new_h = min(needed_in(bullets), bot_lim - top_lim)
        new_oy = max(top_lim, min(center - new_h / 2.0, bot_lim - new_h))
        set_geom(shape, ox, new_oy * EMU, cx, new_h * EMU)
        cy = int(new_h * EMU)
        fit_note = f"grew h->{new_h:.2f}in"
        while (not staggered) and needed_in(bullets) > cy / EMU + 1e-6 and len(bullets) > 1:
            bullets.pop()
            fit_note += f", dropped to {len(bullets)}"

    # --- write bullets (one paragraph each) + hanging indent ------------------
    paras = txBody.findall("a:p", NS)
    anim_note = "unchanged (staggered)"
    if staggered:                                   # fill existing animated paras
        for i, slot in enumerate(animated):
            p = paras[slot]
            rpr = _para_rpr(p, txBody)
            _set_pPr_hanging(p, marL_emu)
            _rebuild_para(p, [bul(bullets[i])], rpr)
    else:                                           # rebuild box with N bullets
        # Take run/paragraph formatting from the first run-bearing paragraph, then
        # DROP every existing paragraph (including any template placeholder bullets)
        # so only our bullets remain.
        tmpl = next((p for p in paras if p.find("a:r", NS) is not None),
                    paras[0] if paras else None)
        if tmpl is not None:
            rpr = _para_rpr(tmpl, txBody)
            ppr_base = tmpl.find("a:pPr", NS)
            ppr_base = copy.deepcopy(ppr_base) if ppr_base is not None else None
        else:
            rpr, ppr_base = _template_rpr(txBody), None
        for p in paras:
            txBody.remove(p)
        for b in bullets:
            txBody.append(make_bullet_para(txBody, rpr, ppr_base, marL_emu, bul(b)))
        if not bullets:                             # keep a valid, empty txBody
            txBody.append(txBody.makeelement(qn("a:p"), {}))
        if animated:
            anim_note = update_animation_range(slide, spid, 0, len(bullets))

    return bullets, fit_note, anim_note


# ---------------------------------------------------------------------------
# Titles: widen to one line, avoid content overlap (font size never changes)
# ---------------------------------------------------------------------------
def fit_title_width(slide, typeface, pt, slide_w_in):
    """Widen the AI_title box horizontally so its (single-line) heading fits
    without wrapping to a second row. Grows to the RIGHT, shifting left only if
    the box would run off the slide; never shrinks the box or the font. Height
    and vertical position are left for fix_title_geometry.

    Returns {'widened_in', 'box_w_in', 'need_in', 'fit'} or None if no AI_title.
    """
    title = find_shape(slide, "AI_title")
    if title is None:
        return None
    text = title.text_frame.text.replace("\x0b", " ").strip()
    ins = get_ins(title)
    ox, oy, cx, cy = geom_emu(title)
    need_usable = text_width_in(text, typeface, pt) + TITLE_WIDTH_SLACK_IN
    cur_usable = cx / EMU - ins["l"] - ins["r"]
    if need_usable <= cur_usable:
        return {"widened_in": 0.0, "box_w_in": cx / EMU,
                "need_in": need_usable, "fit": True}

    margin = MIN_TOP_IN
    max_w = slide_w_in - 2 * margin
    new_cx_in = min(need_usable + ins["l"] + ins["r"], max_w)
    ox_in = ox / EMU
    if ox_in + new_cx_in > slide_w_in - margin:          # would overflow right
        ox_in = max(margin, slide_w_in - margin - new_cx_in)
    set_geom(title, ox_in * EMU, oy, new_cx_in * EMU, cy)
    fit = (new_cx_in - ins["l"] - ins["r"]) + 1e-6 >= need_usable
    return {"widened_in": new_cx_in - cx / EMU, "box_w_in": new_cx_in,
            "need_in": need_usable, "fit": fit}


def fix_title_geometry(slide, typeface, pt):
    """Set the AI_title box to a one-line height and nudge it up if it would
    overlap the AI_content box. Width and left position are never changed."""
    title = find_shape(slide, "AI_title")
    content = find_shape(slide, "AI_content")
    if title is None or content is None:
        return None
    ins = get_ins(title)
    one_line = LINE_FACTOR * pt / 72.0 + ins["t"] + ins["b"]      # inches
    tox, toy, tcx, tcy = geom_emu(title)
    _, coy, _, _ = geom_emu(content)
    new_tcy = one_line * EMU
    new_toy = toy
    title_bottom = toy + new_tcy
    content_top = coy
    moved = 0.0
    if title_bottom > content_top - GAP_IN * EMU:                 # overlap -> up
        new_toy = content_top - GAP_IN * EMU - new_tcy
        new_toy = max(new_toy, MIN_TOP_IN * EMU)
        moved = (toy - new_toy) / EMU
    set_geom(title, tox, new_toy, tcx, new_tcy)
    return {"one_line_in": one_line, "moved_up_in": moved,
            "gap_in": (coy - (new_toy + new_tcy)) / EMU}


# ---------------------------------------------------------------------------
# Unique names (prevent morphing of the text boxes) -- used by the dummy-fill
# test path (run_build_tests); the AI pipeline leaves names shared on purpose
# so Bar/Hex slides morph into each other.
# ---------------------------------------------------------------------------
def rename_unique(prs):
    n = 0
    for sidx, slide in enumerate(prs.slides, start=1):
        for sh in iter_all_shapes(slide.shapes):
            if sh.name == "AI_content":
                sh.name = f"AI_content_{sidx}"; n += 1
            elif sh.name == "AI_title":
                sh.name = f"AI_title_{sidx}"; n += 1
            # AI_tag_*, AI_subtitle, USER_info, decorations: left as-is
    return n


def delete_slides(prs, indices_to_remove, verbose=True):
    sldIdLst = prs.slides._sldIdLst
    sldIds = list(sldIdLst)
    for idx in sorted(indices_to_remove, reverse=True):
        if not 0 <= idx < len(sldIds):
            print(f"  delete_slides: index {idx} out of range, skipping")
            continue
        rId = sldIds[idx].get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldIds[idx])
        if verbose:
            print(f"  deleted slide at index {idx} (rId={rId})")
    return prs


# ---------------------------------------------------------------------------
# Spec loading
# ---------------------------------------------------------------------------
def load_spec(template_name):
    """Load templates/<template_name>/spec.json. Adds '_dir' and
    '_template_path' (the resolved template.pptx path) for convenience."""
    folder = os.path.dirname(os.path.abspath(__file__))
    spec_dir = os.path.join(folder, "templates", template_name)
    spec_path = os.path.join(spec_dir, "spec.json")
    if not os.path.exists(spec_path):
        raise FileNotFoundError(f"No spec.json for template {template_name!r} "
                                 f"(expected {spec_path})")
    with open(spec_path, encoding="utf-8") as f:
        spec = json.load(f)
    spec["_dir"] = spec_dir
    spec["_template_path"] = os.path.join(spec_dir, spec["file"])
    if not os.path.exists(spec["_template_path"]):
        raise FileNotFoundError(f"Template file not found: {spec['_template_path']}")
    return spec


# ---------------------------------------------------------------------------
# Slide-count selection rules (Start + Wheel + Bars + Hexagon + End)
# ---------------------------------------------------------------------------
def _extra_block_size(remaining, min_size, max_size):
    """Size of the next extra-section block (looping phase beyond total_slides)."""
    if remaining >= min_size + max_size:
        return max_size
    if remaining >= min_size:
        return min_size
    return remaining          # < min_size: take what's left


def _fill_extra(extra, n_bars, n_hex, block_size_range):
    """Extra slides beyond spec['total_slides'], inserted between the last
    Hexagon slide and End. Alternating Bars/Hexagon blocks starting with Bars;
    bars and hexagons each advance through their own continuous counter,
    looping back through the template's Bar/Hexagon slides."""
    min_size, max_size = block_size_range
    labels = []
    remaining, want, b, h = extra, "B", 0, 0
    while remaining > 0:
        size = _extra_block_size(remaining, min_size, max_size)
        for _ in range(size):
            if want == "B":
                labels.append(f"B{b % n_bars + 1}"); b += 1
            else:
                labels.append(f"H{h % n_hex + 1}"); h += 1
        remaining -= size
        want = "H" if want == "B" else "B"
    return labels


def build_slide_sequence(count, spec):
    """Ordered variant sequence for a `count`-slide deck, driven entirely by
    `spec` (sections, slide_rules, loop_rules). Returns (kind, template_index,
    label) tuples. Raises ValueError if `count` is out of the spec's range."""
    min_slides = spec.get("min_slides", 6)
    max_slides = spec.get("max_slides", 30)
    if count < min_slides:
        raise ValueError(f"Minimum {min_slides} slides required.")
    if count > max_slides:
        raise ValueError(f"Maximum {max_slides} slides allowed.")

    sections = spec["sections"]
    start_idx = sections["start"]["slide"]
    end_idx = sections["end"]["slide"]
    wheel_idx, n_wheel = sections["wheel"]["slides"], sections["wheel"]["count"]
    bar_idx, n_bars = sections["bars"]["slides"], sections["bars"]["count"]
    hex_idx, n_hex = sections["hexagon"]["slides"], sections["hexagon"]["count"]
    total_slides = spec["total_slides"]

    labels = ["S"] + [f"W{i + 1}" for i in range(n_wheel)]
    if count <= total_slides:
        rule = spec["slide_rules"].get(str(count))
        if rule is None:
            raise ValueError(f"No slide_rules entry for count={count}")
        labels += [f"B{i + 1}" for i in range(rule["bars"])]
        labels += [f"H{i + 1}" for i in range(rule["hex"])]
    else:                                            # beyond total_slides: loop extras
        labels += [f"B{i + 1}" for i in range(n_bars)]
        labels += [f"H{i + 1}" for i in range(n_hex)]
        labels += _fill_extra(count - total_slides, n_bars, n_hex,
                               spec["loop_rules"]["block_size_range"])
    labels += ["E"]

    idx = {"S": start_idx, "E": end_idx}
    for i in range(n_wheel):
        idx[f"W{i + 1}"] = wheel_idx[i]
    for i in range(n_bars):
        idx[f"B{i + 1}"] = bar_idx[i]
    for i in range(n_hex):
        idx[f"H{i + 1}"] = hex_idx[i]
    kind = {"S": "Start", "W": "Wheel", "B": "Bars", "H": "Hexagon", "E": "End"}

    seq = [(kind[lab[0]], idx[lab], lab) for lab in labels]
    assert len(seq) == count, f"sequence length {len(seq)} != {count}"
    return seq


def describe_sequence(count, spec):
    """Print the full labeled variant sequence, e.g. '20 slides: S W1 ... E'."""
    seq = build_slide_sequence(count, spec)
    print(f"{count} slides: {' '.join(t[2] for t in seq)}")
    return seq


# ---------------------------------------------------------------------------
# Slide cloning + deck assembly (loops bars/hex by duplicating template slides)
# ---------------------------------------------------------------------------
def clone_slide(prs, src_slide):
    """Append a full copy of `src_slide` to `prs`. Shares media parts, remaps
    relationship ids in the copied XML, drops the per-slide notes rel, and keeps
    the morph transition. Returns the new SlidePart."""
    package = prs.part.package
    src = src_slide.part
    partname = package.next_partname("/ppt/slides/slide%d.xml")
    new_el = copy.deepcopy(src._element)
    new_part = SlidePart(partname, src.content_type, package, new_el)

    rid_map = {}
    for rId, rel in src.rels.items():
        if rel.reltype.endswith("notesSlide"):
            continue                            # notes belong to one slide
        if rel.is_external:
            new_rId = new_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_rId = new_part.relate_to(rel.target_part, rel.reltype)
        rid_map[rId] = new_rId
    for el in new_el.iter():                    # point copied r:embed/r:id at new rIds
        for attr, val in list(el.attrib.items()):
            if attr.startswith(R_NS) and val in rid_map:
                el.set(attr, rid_map[val])

    rId = prs.part.relate_to(new_part, RT.SLIDE)
    prs.slides._sldIdLst.add_sldId(rId)
    return new_part


def open_template(src_path):
    """Open the template, hydrating a OneDrive cloud-only placeholder if needed.

    A dehydrated OneDrive file is a reparse point that python-pptx can't open
    (PackageNotFoundError) and that Python's own read can't copy (PermissionError).
    The Windows shell copy engine (Copy-Item) does trigger the on-demand download,
    so fall back to that and open the resulting local copy.
    """
    try:
        return Presentation(src_path)
    except (PackageNotFoundError, PermissionError):
        import subprocess, tempfile
        local = os.path.join(tempfile.gettempdir(), "slidemorph_template_local.pptx")
        subprocess.run(
            ["powershell", "-NoProfile", "-Command",
             f"Copy-Item -LiteralPath \"{src_path}\" -Destination \"{local}\" -Force"],
            check=True, capture_output=True,
        )
        print(f"  [note] template was cloud-only; hydrated to local copy: {local}")
        return Presentation(local)


def build_deck(count, src_path, out_path, spec):
    """Build a `count`-slide deck from the template per build_slide_sequence and
    save it to out_path. Returns the resulting slide count."""
    seq = build_slide_sequence(count, spec)
    prs = open_template(src_path)
    n_orig = len(prs.slides._sldIdLst)
    for _variant, tmpl_idx, _label in seq:      # clone in final order
        clone_slide(prs, prs.slides[tmpl_idx - 1])
    delete_slides(prs, list(range(n_orig)), verbose=False)  # drop the originals
    prs.save(out_path)
    return len(prs.slides._sldIdLst)


def compress_images(prs, max_dim_px=1600, jpeg_quality=80):
    """Placeholder -- not implemented yet."""
    return prs


# ---------------------------------------------------------------------------
# Reporting
# ---------------------------------------------------------------------------
def print_inserted(prs):
    print("\n" + "=" * 62)
    print("FULL TEXT INSERTED (by slide / shape)")
    print("=" * 62)
    for sidx, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            nm = sh.name
            if not nm.startswith("AI_"):
                continue
            if not sh.has_text_frame:
                continue
            raw = sh.text_frame.text.replace("\x0b", "\n")
            lines = [ln for ln in raw.split("\n") if ln.strip()]
            print(f"\n[Slide {sidx}] {nm}")
            for ln in lines:
                print(f"    {ln}")
    ui = find_shape(prs.slides[0], "USER_info")
    ui_txt = ui.text_frame.text.replace("\n", " | ") if ui else "(missing)"
    print(f"\n[Slide 1] USER_info  ->  UNTOUCHED: {ui_txt!r}")


def cleanup_numbered_outputs(folder):
    """Delete previously generated output_<number>.pptx files (leaves the
    output_test_*.pptx / output_<template>.pptx deliverables alone)."""
    removed = []
    for f in os.listdir(folder):
        stem = f[len("output_"):-len(".pptx")] if (
            f.startswith("output_") and f.endswith(".pptx")) else ""
        if stem.isdigit():
            os.remove(os.path.join(folder, f))
            removed.append(f)
    return sorted(removed)


def run_build_tests(template_name="neon"):
    """Print compositions/sequences and build output_[count].pptx for a set of
    counts; confirm out-of-range counts are rejected."""
    folder = os.path.dirname(os.path.abspath(__file__))
    spec = load_spec(template_name)
    src = spec["_template_path"]
    max_slides = spec.get("max_slides", 30)
    min_slides = spec.get("min_slides", 6)
    valid = [6, 9, 13, 17, 20, 23, 25, 27, 30]
    valid = [c for c in valid if min_slides <= c <= max_slides]
    reject = [min_slides - 1, max_slides + 1]

    print("=" * 64)
    print("FULL SEQUENCES")
    print("=" * 64)
    for c in valid:
        describe_sequence(c, spec)
        print()

    print("=" * 64)
    print("RANGE VALIDATION")
    print("=" * 64)
    for c in reject:
        try:
            build_slide_sequence(c, spec)
            print(f"  {c}: ERROR - was NOT rejected!")
        except ValueError as e:
            print(f"  {c}: rejected -> \"{e}\"")

    print("\n" + "=" * 64)
    print("BUILDING output_[count].pptx")
    print("=" * 64)
    removed = cleanup_numbered_outputs(folder)
    print(f"Cleaned {len(removed)} old file(s): {', '.join(removed) or '(none)'}")
    print(f"Source: {src}")
    for c in valid:
        out = os.path.join(folder, f"output_{c}.pptx")
        try:
            n = build_deck(c, src, out, spec)
            status = "OK" if n == c else f"MISMATCH (got {n})"
            print(f"  output_{c}.pptx  -> {n} slides  [{status}]")
        except Exception as e:  # noqa: BLE001 - report and continue
            print(f"  output_{c}.pptx  -> ERROR: {type(e).__name__}: {e}")


# ---------------------------------------------------------------------------
# Gemini AI content generation
# ---------------------------------------------------------------------------
STAGE1_PROMPT = """You are creating a presentation outline. Given a topic and a slide count, plan the deck.

Topic: {TOPIC}
Total content slides: {content_slide_count} (this is TARGET_SLIDE_COUNT minus 2, since Start and End are not content)

Build a NATURAL NARRATIVE FLOW that fits THIS specific topic. Every presentation should:

1. START WITH INTRODUCTION / CONTEXT -- what is this topic, why does it matter, what
   will be covered. Give the audience footing before any technical content.
2. PROGRESS FROM GENERAL TO SPECIFIC -- introduce the big picture first, then drill into
   details, examples, and specifics. Never open with a technical detail before
   introducing what it belongs to.
3. GROUP RELATED CONTENT TOGETHER -- if the topic has categories or types, discuss each
   category as a coherent block before moving to the next. Do not jump between unrelated
   details.
4. GIVE EACH SLIDE ONE FOCUS -- the slide heading defines exactly what that slide
   discusses, and the content stays on that one focus. Do not cram multiple topics onto
   one slide.
5. END WITH SYNTHESIS -- applications, key takeaways, or conclusions that tie things
   together.
6. MAKE SECTIONS CONNECT -- each new section should logically follow from the previous
   one, so the reader feels a story unfolding, not a random list of facts.

Add examples where they help -- a formula, an equation, a syntax/code snippet, a short
calculation, or a concrete case -- whenever the point calls for it. ANY subject may need
these: business has ROI and break-even formulas, biology has growth equations, economics
has calculations, programming has code/syntax. Do NOT assume the subject has or lacks
technical content; that is decided per point in the next stage. Here, just plan headings
and per-slide topics so the flow is coherent.

The specific sections depend ENTIRELY on the topic -- do NOT force a template. For a
technical topic you might go concepts -> categories -> examples -> applications. For a
historical topic: context -> timeline -> key figures -> consequences. For a business
topic: problem -> solution -> market -> execution. Choose whatever flow makes THIS topic
understandable.

The slides are grouped by PRESENTATION ROLE (structure, not subject). Use these roles to
place your narrative:
- WHEEL slides (the first 3 content slides): introduce and DEFINE the main
  concepts/pillars of the topic. Each wheel slide focuses on ONE pillar. The three
  wheel_tags ARE those three pillars, in order -- wheel_tags[0] is the first wheel
  slide's focus, wheel_tags[1] the second, wheel_tags[2] the third. Make the three wheel
  sections define exactly those three pillars, one per slide, in the same order.
- BAR slides (the middle): the main body -- detailed explanations, methods, categories,
  worked examples. One subtopic per slide, grouped so related bars sit together.
- HEXAGON slides (later): applications, use cases, comparisons, evaluations, or the
  closing synthesis.

For this presentation there are:
- 3 Wheel slides
- {bar_count} Bar slides
- {hex_count} Hexagon slides

For sections with detailed content -- full code examples, step-by-step calculations,
algorithms, detailed explanations, or important quotes -- mark them as
"needs_long_content": true to get a large dedicated slide. Fill these slides with
substantial content (8-18 lines for code, 6-10 sentences for paragraphs) -- these are big
slides that need real content, not short summaries. For decks over 20 slides, always
include at least one long-content section. For shorter decks, use when the topic
genuinely has detailed content worth showing. Never mark wheel/intro sections as
long-content.

{long_content_guidance}

A long-content section replaces one Bar or Hexagon slot -- it does not add a slide, and
it never replaces a Wheel slide. If a point is short enough to work as a few bullets,
leave needs_long_content false and let it be a normal Bar/Hexagon slide.

Return ONLY a JSON object, no markdown, no explanation:
{
  "title": "short title, max 20 chars",
  "subtitle": "subtitle, max 30 chars",
  "wheel_tags": ["pillar1", "pillar2", "pillar3"],
  "sections": [
    {"type": "wheel", "variant": 1, "heading": "max 14 chars", "topic": "define ONLY this pillar", "needs_long_content": false},
    {"type": "bar", "variant": 1, "heading": "max 25 chars", "topic": "what this ONE slide covers", "needs_long_content": false},
    {"type": "hexagon", "variant": 1, "heading": "max 25 chars", "topic": "what this ONE slide covers", "needs_long_content": false}
  ]
}

Rules:
- title max 20 characters
- subtitle max 30 characters
- bar and hexagon headings max 25 characters (must fit on ONE line)
- wheel_tags are the 3 pillars shown on the wheel; each may wrap to 2 lines, so up to
  about 36 characters total (max 18 per line). Prefer short labels, but full words
  are fine -- they will wrap rather than being cut off
- The 3 wheel sections define wheel_tags[0], [1], [2] respectively -- one pillar each,
  same order, and their "topic" must stay on that single pillar only.
- Order the bar and hexagon sections so the narrative flows: general before specific,
  related subtopics grouped together, each section a natural continuation of the last.
- Every section "heading" is that slide's single focus; its "topic" must describe only
  what that one slide covers, with no overlap into other slides.
- sections must have exactly {content_slide_count} entries
- Keep all headings short enough to never wrap to a second line"""

STAGE2_PROMPT = """Write the on-slide content for ONE presentation slide.

THIS SLIDE'S FOCUS (discuss ONLY this): {focus}
Topic of this slide: {section.topic}
Heading: {section.heading}
Slide role: {role_note}

Continuity (for context only -- do NOT write about these):
- The previous slide covered: {prev_topic}. Do not repeat it.
- The next slide will cover: {next_topic}. Do not steal that content.

Stay STRICTLY on this slide's focus/heading. Write only about "{focus}". Do not discuss
what other slides cover, and do not turn this into an overview of the whole topic.

For each point, choose the clearest way to present it. Use bullets for descriptive
information. If a specific point is best explained with a formula, equation, code
snippet, calculation, or short example, include that directly instead of just describing
it. This applies to ANY subject -- business, biology, economics, programming, chemistry,
etc. Do not assume a subject has or lacks technical content; decide from the actual point
being made.

The box is the ONLY hard constraint:
- Each line must fit on ONE line in the box: at most about {max_chars} characters.
- The box holds about {capacity} lines. FILL {target_min}-{target_max} of them
  ({fill_lo_pct}-{fill_hi_pct}% of the box). Do not leave the box nearly empty, and do not overflow it.
- A formula, code line, or calculation counts toward the line budget exactly like a
  bullet does (one line each, unless it is short enough to share).
- Never drop a major point. If space is tight, SHORTEN the wording of a point rather
  than removing the point.{extra}

Style:
- Descriptive lines: clean concrete phrases, start with a capital, no trailing period.
- Formulas / equations / code / calculations: write them in their natural form --
  symbols, operators, and numbers are expected (e.g. "ROI = (Gain - Cost) / Cost * 100").

Return ONLY a JSON object, no markdown, no explanation:
{
  "lines": ["line 1", "line 2", "line 3"]
}"""

STAGE2_LONG_PROMPT = """Write the FULL-DETAIL content for ONE large presentation slide.

Topic of this slide: {section.topic}
Heading: {section.heading}

Stay STRICTLY on this slide's heading/topic. Do not drift into material that belongs on
other slides; this one block is about "{section.heading}" only.

This slide has a MASSIVE box (10.9 inches wide, 5 inches tall). It exists because the
point needs a big block that does NOT work as short bullets -- a full code example (a
whole loop/function/program), a step-by-step calculation with the working shown, a
complete algorithm or pseudocode, a detailed explanatory paragraph, or an important quote
plus its context. Produce exactly that: the real thing, as it should literally appear on
the slide.

How much content (this is the important part -- do NOT under-fill):
- Code blocks: 8-18 lines of actual code with proper indentation.
- Paragraphs: 6-10 sentences of detailed explanation, in full flowing prose.
- Calculations: every step on its own line, the full working shown start to finish.
- Algorithms: the complete pseudocode or numbered steps written out.
- Quotes: the full quote, then a few sentences of context explanation.

Hard constraints:
- Write AT LEAST {min_lines} lines and at most about {max_lines} lines; each line at most
  about {max_chars} characters.
- Three short lines in a huge box is a FAILURE. Use the space.
- Preserve meaningful line breaks AND indentation. For code, put each statement on its
  own line and keep indentation using real spaces (e.g. 4 spaces per level).
- Do NOT use bullet points, and do NOT prefix lines with "-", "*", or "|". Write it
  exactly as it should read.
- For paragraphs, write full sentences that flow; do not compress into fragments.

Return ONLY a JSON object, no markdown, no code fences. Put the whole block in one
string with "\\n" between lines (indentation as literal spaces):
{
  "content": "line one\\nline two\\n    indented line three"
}"""


def _read_env_value(key, env_path):
    """Read a single KEY=value from a .env file (value is never printed)."""
    if not os.path.exists(env_path):
        return None
    for line in open(env_path, encoding="utf-8"):
        line = line.strip()
        if line.startswith("#") or "=" not in line:
            continue
        k, v = line.split("=", 1)
        if k.strip() == key:
            return v.strip().strip('"').strip("'")
    return None


def _is_quota_error(exc):
    """True if `exc` looks like a Gemini quota / rate-limit (429) failure."""
    s = f"{type(exc).__name__}: {exc}".lower()
    return ("resourceexhausted" in s or "429" in s or "quota" in s
            or "exhausted" in s or "rate limit" in s)


def _is_unavailable_error(exc):
    """True if `exc` means the model itself is gone/unusable (404, retired,
    "no longer available", "not found") -- a reason to skip to the next model."""
    s = f"{type(exc).__name__}: {exc}".lower()
    return ("notfound" in s or "404" in s or "not found" in s
            or "no longer available" in s or "not supported" in s)


def _retry_seconds(exc):
    """Seconds the API suggests waiting before retrying (from a 429), or None."""
    s = f"{exc}"
    m = re.search(r"retry in ([\d.]+)s", s) or re.search(r"seconds:\s*(\d+)", s)
    return float(m.group(1)) if m else None


class _FallbackModel:
    """Wraps a chain of Gemini models. On a quota/429 error it switches to the
    next model in the chain (for this call and every later one) and retries,
    so a mid-run quota exhaustion never loses progress. Any non-quota error, or
    running out of fallbacks, propagates normally."""

    def __init__(self, genai, model_names):
        self._genai = genai
        self._names = list(model_names)
        self._i = 0
        self._model = genai.GenerativeModel(self._names[0])
        self._lock = threading.Lock()

    @property
    def name(self):
        return self._names[self._i]

    def generate_content(self, prompt):
        waits = 0
        while True:
            with self._lock:
                model = self._model
                model_idx = self._i
            try:
                return model.generate_content(prompt)
            except Exception as exc:  # noqa: BLE001 - inspected, then re-raised
                quota_delay = None
                with self._lock:
                    skip = _is_quota_error(exc) or _is_unavailable_error(exc)
                    if skip and self._i != model_idx:
                        waits = 0
                    elif skip and self._i + 1 < len(self._names):
                        old = self._names[self._i]
                        reason = "quota" if _is_quota_error(exc) else "unavailable"
                        self._i += 1
                        print(f"    [{reason}] {old} -> falling back to "
                              f"{self._names[self._i]}")
                        self._model = self._genai.GenerativeModel(self._names[self._i])
                        waits = 0
                    elif _is_quota_error(exc) and waits < MAX_QUOTA_WAITS:
                        quota_delay = min(_retry_seconds(exc) or QUOTA_RETRY_CAP,
                                          QUOTA_RETRY_CAP)
                        waits += 1
                        print(f"    [quota] {self.name} rate-limited; waiting "
                              f"{quota_delay:.0f}s then retrying "
                              f"({waits}/{MAX_QUOTA_WAITS})")
                    else:
                        raise
                time.sleep(quota_delay if quota_delay is not None
                           else API_CALL_DELAY)


def load_gemini(folder):
    """Configure Gemini from GEMINI_API_KEY (env var first, then .env);
    return a model handle."""
    import google.generativeai as genai  # lazy: keeps the deprecation warning out of other runs
    key = os.getenv("GEMINI_API_KEY") or _read_env_value(
        "GEMINI_API_KEY", os.path.join(folder, ".env"))
    if not key:
        raise RuntimeError("GEMINI_API_KEY not found in env or .env")
    genai.configure(api_key=key)
    print(f"Gemini ready (models={GEMINI_MODELS}, key loaded: {len(key)} chars)")
    return _FallbackModel(genai, GEMINI_MODELS)


def _resp_text(resp):
    """Extract text from a Gemini response, tolerating blocked/empty results."""
    try:
        if resp.text:
            return resp.text
    except Exception:
        pass
    parts = []
    for c in getattr(resp, "candidates", None) or []:
        for p in getattr(getattr(c, "content", None), "parts", None) or []:
            parts.append(getattr(p, "text", "") or "")
    return "".join(parts)


def _strip_fences(text):
    """Drop ```json ... ``` markdown fences Gemini sometimes adds."""
    t = text.strip()
    if t.startswith("```"):
        nl = t.find("\n")
        t = t[nl + 1:] if nl != -1 else t[3:]
        if "```" in t:
            t = t[:t.rfind("```")]
    return t.strip()


def gemini_json(model, prompt, retries=1):
    """Call Gemini, strip fences, parse JSON. Retry once on invalid JSON."""
    last = None
    for attempt in range(retries + 1):
        t0 = time.time()
        raw = _resp_text(model.generate_content(prompt))
        dt = time.time() - t0
        print(f"    [perf] Gemini API call took {dt:.2f}s (model={model.name})")
        try:
            return json.loads(_strip_fences(raw))
        except (json.JSONDecodeError, ValueError) as e:
            last = e
            if attempt < retries:
                print("    [warn] invalid JSON from Gemini, retrying once...")
                time.sleep(API_CALL_DELAY)
    raise ValueError(f"invalid JSON after retry: {last}")


def stage2_bullets(model, section, max_chars, capacity, target_min, target_max, rn,
                   fill_lo_pct=70, fill_hi_pct=85, extra="",
                   focus="", prev_topic="", next_topic=""):
    """Ask Gemini for the slide's content lines. Each returned line is whatever
    best communicates that point -- a bullet, a formula, a code line, or a short
    calculation -- constrained only by the box (max_chars per line, ~capacity
    lines, fill target_min..target_max, i.e. fill_lo_pct-fill_hi_pct% of the box).
    `extra` carries optional per-section guidance (empty for box-only sections).
    `focus` is the single thing this slide must discuss (the middle wheel tag, or
    the bar/hexagon heading); `prev_topic`/`next_topic` give continuity so the
    slide neither repeats the previous slide nor steals the next one's content.
    Returns a list of raw line strings."""
    focus = focus or section.get("focus_tag") or section.get("heading", "")
    prev_topic = (prev_topic or "").strip() or "(nothing -- this is the first content slide)"
    next_topic = (next_topic or "").strip() or "(nothing -- this is the last content slide)"
    prompt = (STAGE2_PROMPT
              .replace("{focus}", str(focus))
              .replace("{prev_topic}", prev_topic)
              .replace("{next_topic}", next_topic)
              .replace("{section.topic}", str(section.get("topic", "")))
              .replace("{section.heading}", str(section.get("heading", "")))
              .replace("{role_note}", str(rn))
              .replace("{max_chars}", str(max_chars))
              .replace("{capacity}", str(capacity))
              .replace("{target_min}", str(target_min))
              .replace("{target_max}", str(target_max))
              .replace("{fill_lo_pct}", str(fill_lo_pct))
              .replace("{fill_hi_pct}", str(fill_hi_pct))
              .replace("{extra}", str(extra)))
    data = gemini_json(model, prompt)
    if not isinstance(data, dict):
        return []
    return data.get("lines") or data.get("bullets") or []


def stage2_long_content(model, section, max_lines, max_chars, min_lines=1):
    """Ask Gemini for the full multi-line block for a long-content slide (code,
    derivation, algorithm, or paragraph). Returns the block as a single string
    with '\\n' line breaks and literal-space indentation preserved."""
    prompt = (STAGE2_LONG_PROMPT
              .replace("{section.topic}", str(section.get("topic", "")))
              .replace("{section.heading}", str(section.get("heading", "")))
              .replace("{min_lines}", str(min_lines))
              .replace("{max_lines}", str(max_lines))
              .replace("{max_chars}", str(max_chars)))
    data = gemini_json(model, prompt)
    if isinstance(data, dict):
        body = data.get("content")
        if isinstance(body, list):                 # tolerate a list of lines
            body = "\n".join(str(x) for x in body)
        if isinstance(body, str) and body.strip():
            return body
    return ""


def truncate(text, n, ellipsis=False):
    """Clip `text` to `n` chars. With ellipsis=True an over-limit string is cut to
    n-3 chars and gets a trailing '...' so the RESULT is still <= n chars."""
    text = (text or "").strip()
    if len(text) <= n:
        return text
    if ellipsis and n > 3:
        return text[:n - 3].rstrip() + "..."
    return text[:n].rstrip()


# ---------------------------------------------------------------------------
# Adaptive content generation, driven ONLY by box size (Pillow-measured).
#
# No subject assumptions and no per-type "bullets-only" rule: each returned line
# is whatever best communicates its point (bullet, formula, code, calculation),
# and the sole hard constraint is the box -- ~max_chars per line and a line
# capacity computed from the box height. We aim to FILL 70-85% of the box by
# default; a section's spec may override this (fill_target/max_lines) to keep a
# smaller box -- e.g. the wheel intro boxes -- less packed.
# ---------------------------------------------------------------------------
CONTENT_LINE_H_IN = LINE_FACTOR * CONTENT_PT / 72.0   # one content line's height
FILL_TARGET_LO = 0.70      # default: aim to fill at least this fraction of the box
FILL_TARGET_HI = 0.85      # ... and at most this fraction
FILL_UNDERFULL = 0.60      # below this -> regenerate asking for more detail

_CHARS_SAMPLE = ("The quick brown fox jumps over the lazy dog while counting "
                 "1234567890 items and notes")


def count_wraps(bullets, usable_in, hang_in):
    """How many `bullets` are too wide to fit on ONE line in the content box
    (measured with Pillow via count_display_lines, including the '- ' hang)."""
    return sum(1 for b in bullets
               if count_display_lines(bul(b), usable_in, hang_in) > 1)


def total_display_lines(items, usable_in, hang_in):
    """Total rendered line count for `items` (each may wrap), measured w/ Pillow."""
    return sum(count_display_lines(bul(x), usable_in, hang_in) for x in items)


def approx_chars_per_line(usable_in):
    """Roughly how many content-font characters fit on one usable-width line.
    Measured with Pillow against a mixed-case sample (real gate stays the
    per-line wrap measurement; this is just guidance passed to the model)."""
    avg = text_width_in(_CHARS_SAMPLE, "+mj-lt", CONTENT_PT) / len(_CHARS_SAMPLE)
    return max(10, int(usable_in / avg))


def box_line_capacity(cy_emu, ins):
    """How many content lines fit in a box of height `cy_emu` (EMU) given insets."""
    usable_h = cy_emu / EMU - ins["t"] - ins["b"]
    return max(1, int(usable_h / CONTENT_LINE_H_IN))


def wheel_focus_index(slide):
    """Which AI_tag_N shape sits in the wheel's MIDDLE (focus) position on this
    wheel slide. The template stacks two tags on the left and places the third
    near the wheel (furthest to the right); that right-most tag is the slide's
    focus. Returns the 1-based tag index N (so its text is wheel_tags[N-1]), or
    None if the slide has no AI_tag_* shapes. Purely geometry-driven, so it works
    for any wheel layout without hardcoding which slide focuses on which tag."""
    best_n, best_cx = None, None
    for n in (1, 2, 3):
        sh = find_shape(slide, f"AI_tag_{n}")
        if sh is None:
            continue
        ox, _oy, cx, _cy = geom_emu(sh)
        center_x = ox + cx / 2.0
        if best_cx is None or center_x > best_cx:
            best_cx, best_n = center_x, n
    return best_n


def wheel_focus_tag(slide, wheel_tags):
    """The text of this wheel slide's focus (middle) tag, from wheel_tags. Falls
    back to the first tag if geometry detection or the index is unavailable."""
    n = wheel_focus_index(slide)
    if n is not None and wheel_tags and 1 <= n <= len(wheel_tags):
        return wheel_tags[n - 1]
    return wheel_tags[0] if wheel_tags else ""


def role_note(stype):
    """Structural (NOT subject) role blurb passed to the Stage 2 prompt."""
    return {
        "wheel":   "Wheel slide -- introduce/define a core concept (shorter box).",
        "bar":     "Bar slide -- main body; explain this subtopic in detail.",
        "bars":    "Bar slide -- main body; explain this subtopic in detail.",
        "hexagon": "Hexagon slide -- applications, examples, or summary.",
    }.get((stype or "").lower(), "Content slide -- explain this point clearly.")


def _looks_technical(line):
    """Heuristic marker for the printout: does this line read as a formula,
    equation, calculation, or code (rather than a descriptive phrase)? Used only
    for reporting so a human can verify technical content appeared where needed."""
    s = line.strip()
    has_math = bool(re.search(r"[=×÷≈≤≥∑∏√^]", s)) or \
        bool(re.search(r"\d\s*[-+*/×÷]\s*\d", s))
    has_code = bool(re.search(r"[A-Za-z_]\w*\([^)]*\)", s)) or "()" in s or \
        s.endswith(":") or bool(re.search(r"[{}<>]|::|=>|->", s))
    return has_math or has_code


def _content_extra(spec_cfg, line_cap):
    """Optional per-section guidance appended to the Stage 2 prompt. Only sections
    that declare an explicit `max_bullets` in the spec (currently the wheel intro
    boxes) get it; box-only sections (bars/hexagon) return "" so their prompt is
    unchanged."""
    mb = spec_cfg.get("max_bullets")
    if not mb:
        return ""
    wpb = spec_cfg.get("words_per_bullet", [4, 6])
    wrap = spec_cfg.get("allow_wrap", [1, 2])
    return ("\n- This is an intro/definition box: aim for about "
            f"{mb[0]}-{mb[1]} substantive bullets of roughly {wpb[0]}-{wpb[1]} "
            f"words each, at most {line_cap} lines total. Most bullets should sit "
            f"on one line, but {wrap[0]}-{wrap[1]} may wrap to a second line. Fill "
            "the box without overflowing, and do not drop a key point.")


def adaptive_bullets(model, section, usable_in, hang_in, ph_content_spec, capacity,
                     focus="", prev_topic="", next_topic=""):
    """Generate the slide's content lines, filling a target fraction of the box.

    `focus` is the single thing the slide must discuss (wheel middle tag or
    bar/hexagon heading); `prev_topic`/`next_topic` are the neighbouring slides'
    focuses, passed to Stage 2 for continuity so slides don't overlap.

    Requests content sized to the box (max_chars/line + fill target derived from
    `capacity`), measures the rendered line total with Pillow, and regenerates
    once or twice if the box is under-full or overflowing. Points are never
    dropped here -- the model is asked to shorten instead. Returns
    (lines, total_lines, log).

    The fill target defaults to FILL_TARGET_LO..HI (70-85%), but a section's spec
    (`ph_content_spec`) may override it with `fill_target` [lo, hi] and cap the
    rendered lines with `max_lines` -- used to keep the smaller wheel intro boxes
    (70-80%, <=10 lines, mostly one-line bullets) less packed. With no overrides
    the behaviour is identical to before, so bars/hexagon are unaffected.
    """
    spec_cfg = ph_content_spec or {}
    lo, hi = spec_cfg.get("fill_target", [FILL_TARGET_LO, FILL_TARGET_HI])
    underfull = spec_cfg.get("fill_underfull", FILL_UNDERFULL)
    # Hard ceiling on rendered lines (default: the whole box). Lets a shorter
    # intro box (wheel) stay less packed than its raw height would allow.
    line_cap = min(capacity, int(spec_cfg.get("max_lines", capacity)))

    max_chars = approx_chars_per_line(usable_in)
    target_min = max(1, round(lo * capacity))
    target_max = max(target_min, int(hi * capacity))
    target_max = min(target_max, line_cap)          # never exceed the line cap
    target_min = min(target_min, target_max)
    low_water = min(max(1, int(underfull * capacity)), target_max)
    fill_lo_pct = int(round(lo * 100))
    fill_hi_pct = int(round(hi * 100))
    extra = _content_extra(spec_cfg, line_cap)
    rn = role_note(section.get("type"))
    log = []
    last = []
    for attempt in range(MAX_BULLET_ITERS):
        try:
            raw = stage2_bullets(model, section, max_chars, capacity,
                                 target_min, target_max, rn,
                                 fill_lo_pct, fill_hi_pct, extra,
                                 focus, prev_topic, next_topic)
        except Exception as exc:  # noqa: BLE001 - reported, then fall back
            log.append(f"try{attempt + 1}: gen FAILED: {exc}")
            time.sleep(API_CALL_DELAY)
            break
        time.sleep(API_CALL_DELAY)
        items = [strip_bullet_prefix(x) for x in raw if isinstance(x, str)]
        items = [x for x in items if x]         # drop marker-only / empty lines
        total = total_display_lines(items, usable_in, hang_in)
        pct = (100.0 * total / capacity) if capacity else 0.0
        log.append(f"try{attempt + 1}: {len(items)} items -> {total} lines "
                   f"(cap {capacity}, cap_lines {line_cap}, {pct:.0f}% full; "
                   f"target {target_min}-{target_max})")
        last = items
        if total > line_cap:                    # overflow -> ask to shorten/fewer
            target_max = max(target_min, line_cap - 1)
            continue
        if total < low_water:                   # under-full -> ask for more detail
            target_min = min(line_cap, target_max + 1)
            target_max = line_cap
            continue
        return items, total, log                # in-range -> accept

    # did not converge: keep the last set. fill_content_bullets grows the box to
    # fit, and only drops as an absolute last resort.
    return last, total_display_lines(last, usable_in, hang_in), log


def _fill_shape_text(slide, name, lines):
    sh = find_shape(slide, name)
    if sh is not None and sh.has_text_frame:
        fill_simple(sh, lines)


# ---------------------------------------------------------------------------
# Long-content routing (optional slides 21/22): pair each content slot with its
# outline section, then reroute sections the AI flagged needs_long_content onto
# the template's long-content slides. Structure-driven -- if a template's spec
# has no "long_content" section, routing is simply skipped.
# ---------------------------------------------------------------------------
_KIND_TO_PH = {"Wheel": "wheel", "Bars": "bars", "Hexagon": "hexagon"}


def long_content_limits(count, spec):
    """(min_long, max_long) for a `count`-slide deck, from long_content_rules
    ('by_deck_size' bands, clamped by 'max_per_deck'). Small decks may legitimately
    use zero; long decks (21-30) require at least one. Returns (0, 0) for a
    template with no long-content slides."""
    lc = spec["sections"].get("long_content")
    rules = spec.get("long_content_rules", {})
    if not lc or not lc.get("slides"):
        return 0, 0
    cap = rules.get("max_per_deck", 0)
    lo, hi = 0, cap
    for band in rules.get("by_deck_size", []):
        if band.get("from", 0) <= count <= band.get("to", 10 ** 6):
            lo = band.get("min_long", 0)
            hi = band.get("max_long", cap)
            break
    hi = min(hi, cap)
    return min(lo, hi), hi


def route_plan(sequence, sections, spec):
    """Turn (sequence, outline sections) into an ordered list of slot dicts, one
    per slide: {kind, tmpl_idx, label, section, is_long}. Sections flagged
    needs_long_content (Bars/Hexagon only -- never Wheel) are rerouted onto the
    spec's long-content template slides, alternating between them, within the
    deck-size limits from long_content_limits(). If the outline flagged fewer than
    the required minimum (decks of 21+), extra Bars/Hexagon slots are promoted to
    make up the difference. A rerouted slot REPLACES its bar/hex slot, so the deck
    length never changes. Returns the slot list (len == len(sequence))."""
    lc = spec["sections"].get("long_content")
    rules = spec.get("long_content_rules", {})
    long_slides = lc["slides"] if lc else []
    min_long, max_long = long_content_limits(len(sequence), spec)
    route_from = set(rules.get("route_from", ["bars", "hexagon"]))

    # Assign sections to content slots by type, mirroring the outline's grouping.
    queues = {"wheel": [], "bars": [], "hexagon": []}
    for sec in sections:
        st = (sec.get("type") or "").lower()
        queues.setdefault("bars" if st == "bar" else st, []).append(sec)

    slots = []
    for kind, tmpl_idx, label in sequence:
        if kind in ("Start", "End"):
            slots.append({"kind": kind, "tmpl_idx": tmpl_idx, "label": label,
                          "section": None, "is_long": False})
            continue
        qk = _KIND_TO_PH[kind]
        sec = queues[qk].pop(0) if queues[qk] else {"type": qk, "heading": "",
                                                    "topic": ""}
        slots.append({"kind": kind, "tmpl_idx": tmpl_idx, "label": label,
                      "section": sec, "is_long": False})

    def eligible(slot):
        """Bars/Hexagon only -- Wheel slides are always kept for intro/definitions."""
        return (not slot["is_long"]
                and _KIND_TO_PH.get(slot["kind"]) in route_from)

    def promote(slot, n):
        slot["orig_kind"] = slot["kind"]
        slot["pre_long_tmpl_idx"] = slot["tmpl_idx"]    # so it can be demoted back
        slot["kind"] = "LongContent"
        slot["tmpl_idx"] = long_slides[n % len(long_slides)]
        slot["is_long"] = True

    # Reroute the flagged Bars/Hexagon slots, up to the deck-size maximum.
    used = 0
    for slot in slots:
        if used >= max_long:
            break
        sec = slot["section"]
        if sec and sec.get("needs_long_content") and eligible(slot):
            promote(slot, used)
            used += 1

    # Long decks must have at least min_long. If the outline under-flagged, promote
    # additional Bars/Hexagon slots, preferring the most detail-heavy role (Bars)
    # and skipping the very first content slide.
    if used < min_long:
        pool = [s for s in slots[1:] if eligible(s)]
        pool.sort(key=lambda s: 0 if s["kind"] == "Bars" else 1)
        for slot in pool:
            if used >= min_long:
                break
            slot["forced_long"] = True
            promote(slot, used)
            used += 1

    return slots


def build_deck_from_slots(slots, src_path, out_path):
    """Clone the template slides named by each slot's `tmpl_idx` (in order) into
    a fresh deck, drop the originals, and save. Returns the final slide count."""
    prs = open_template(src_path)
    n_orig = len(prs.slides._sldIdLst)
    for slot in slots:
        clone_slide(prs, prs.slides[slot["tmpl_idx"] - 1])
    delete_slides(prs, list(range(n_orig)), verbose=False)
    prs.save(out_path)
    return len(prs.slides._sldIdLst)


MAX_LONG_ITERS = 3          # generation attempts before accepting an under-full block
LONG_MAX_LINES = 20         # hard cap: never insert more than this many lines
LONG_SAFE_FRAC = 0.85       # fill target ceiling: only use this much of the box height
LONG_GAP_IN = 0.15          # min vertical gap between title bottom and content top


def long_box_spec(spec, shape):
    """Line budget for a long-content AI_content box: (usable_in, capacity,
    max_lines, max_chars, min_lines). `capacity` leaves a safety margin -- only
    LONG_SAFE_FRAC of the box height is used -- and is hard-capped at
    LONG_MAX_LINES, so a filled block can never overflow the box. min_lines is the
    60%-fill floor the model is told to clear."""
    lc_ph = spec["placeholders"]["long_content"]["AI_content"]
    _, _, cx, cy = geom_emu(shape)
    ins = get_ins(shape)
    _tf, pt = run_font(shape, "+mj-lt", CONTENT_PT)
    usable = cx / EMU - ins["l"] - ins["r"]
    usable_h = cy / EMU - ins["t"] - ins["b"]
    line_h = LINE_FACTOR * pt / 72.0
    by_height = int(LONG_SAFE_FRAC * usable_h / line_h)          # 85% of the box
    capacity = max(1, min(by_height, LONG_MAX_LINES))            # ... and <= 20
    max_lines = min(lc_ph.get("max_lines", capacity), capacity)
    # The spec's max_chars_per_line is an upper bound, not the truth: what actually
    # fits is measured from the real box width and font. Taking the larger of the
    # two lets prose lines wrap to 2 display lines each, which silently halves the
    # usable line budget and gets the tail of the block trimmed off.
    measured_chars = approx_chars_per_line(usable)
    max_chars = min(lc_ph.get("max_chars_per_line", measured_chars), measured_chars)
    frac = spec.get("long_content_rules", {}).get("min_fill_fraction", 0.60)
    min_lines = max(1, int(round(frac * capacity)))
    return usable, capacity, max_lines, max_chars, min_lines


def measure_long_body(body, usable_in, capacity):
    """Split a generated block into slide lines and measure it: returns
    (kept_lines, used_display_lines, n_dropped). Leading list markers are stripped
    from every line (long-content is plain text -- no bullets), while code
    indentation on unmarked lines is preserved. Interior blank lines are kept,
    leading and trailing blanks dropped, and the block is trimmed to the box.
    n_dropped > 0 means the block did not fit and was cut -- which leaves the slide
    ending mid-sentence, so callers should regenerate rather than ship it."""
    raw = (body or "").replace("\r\n", "\n").replace("\r", "\n").split("\n")
    raw = [strip_long_bullet(x) for x in raw]           # Problem 4: no bullets here
    while raw and not raw[0].strip():
        raw.pop(0)
    while raw and not raw[-1].strip():
        raw.pop()
    kept, used = [], 0
    for ln in raw:
        dl = count_display_lines(ln, usable_in, 0.0) if ln.strip() else 1
        if kept and used + dl > capacity:
            break
        kept.append(ln)
        used += dl
    return kept, used, len(raw) - len(kept)


def ensure_gap_below_title(slide, min_gap_in):
    """Move the AI_content box down (shrinking its height to match) so its top sits
    at least `min_gap_in` below the AI_title box's bottom edge. No-op when the gap
    is already large enough. Stops the heading from overlapping the content on
    long-content slides. Returns the distance moved, in inches."""
    title = find_shape(slide, "AI_title")
    content = find_shape(slide, "AI_content")
    if title is None or content is None:
        return 0.0
    _, toy, _, tcy = geom_emu(title)
    cox, coy, ccx, ccy = geom_emu(content)
    need_top = toy + tcy + int(round(min_gap_in * EMU))
    if coy >= need_top:
        return 0.0
    delta = need_top - coy
    set_geom(content, cox, need_top, ccx, max(int(0.5 * EMU), ccy - delta))
    return delta / EMU


def trim_to_box_height(shape, lines):
    """Absolute overflow guard: drop lines from the bottom until the block's
    rendered height fits inside 100% of the box's usable height. `capacity` in
    long_box_spec already leaves a margin, so this rarely fires -- it is the final
    'never let text leave the box' backstop. Returns (kept_lines, used_lines)."""
    _, _, cx, cy = geom_emu(shape)
    ins = get_ins(shape)
    usable = cx / EMU - ins["l"] - ins["r"]
    usable_h = cy / EMU - ins["t"] - ins["b"]
    max_disp = max(1, int(usable_h / (LINE_FACTOR * CONTENT_PT / 72.0)))
    kept, used = [], 0
    for ln in lines:
        dl = count_display_lines(ln, usable, 0.0) if ln.strip() else 1
        if kept and used + dl > max_disp:
            break
        kept.append(ln)
        used += dl
    return kept, used


def generate_long_body(model, section, usable_in, capacity, max_lines, max_chars,
                       min_lines):
    """Generate a long-content block that both FILLS the box (>= min_lines display
    lines) and FITS it complete (nothing trimmed). Under-full -> ask for more;
    over-long -> ask for less, since a trimmed block ends mid-sentence or
    mid-statement. Keeps the best candidate seen, preferring complete blocks.
    Returns (kept_lines, used_lines, log)."""
    log = []
    best, best_used, best_score = [], 0, (-1, -1)
    ask_min, ask_max, ask_chars = min_lines, max_lines, max_chars
    for attempt in range(MAX_LONG_ITERS):
        try:
            body = stage2_long_content(model, section, ask_max, ask_chars, ask_min)
        except Exception as exc:  # noqa: BLE001 - reported, then keep what we have
            log.append(f"try{attempt + 1}: gen FAILED: {exc}")
            break
        time.sleep(API_CALL_DELAY)
        kept, used, dropped = measure_long_body(body, usable_in, capacity)
        pct = (100.0 * used / capacity) if capacity else 0.0
        cut = f", CUT {dropped} line(s)" if dropped else ""
        log.append(f"try{attempt + 1}: asked {ask_min}-{ask_max} -> {len(kept)} lines "
                   f"-> {used} display lines ({pct:.0f}% of {capacity}; "
                   f"need >= {min_lines}{cut})")

        score = (0 if dropped else 1, used)      # complete beats fuller-but-cut
        if score > best_score:
            best, best_used, best_score = kept, used, score
        if used >= min_lines and not dropped:
            return kept, used, log

        if dropped:
            # Overflow is usually long lines WRAPPING, not too many lines: prose
            # written at the char limit costs 2 display lines each. Tighten the
            # per-line width as well as the line count.
            wrapped = used > len(kept)
            if wrapped:
                ask_chars = max(40, int(ask_chars * 0.85))
            fitted = max(min_lines, len(kept) - 1)
            ask_max = max(min_lines, min(ask_max - 1, fitted) if not wrapped
                          else ask_max)
            ask_min = min(ask_min, ask_max)
        else:                                    # under-full: push for more
            ask_min = min(capacity, ask_min + max(2, min_lines // 3))
            ask_max = max(ask_max, ask_min)
    return best, best_used, log


def prefetch_long_content(slots, spec, model, src_path, min_long):
    """Generate every long-content block BEFORE the deck is cloned, so a slot whose
    content comes back bullet-sized can be demoted back to a normal Bars/Hexagon
    slide instead of leaving a mostly-empty huge box. Boxes are measured on the
    template's own long-content slides. Caches the accepted block on the slot as
    'long_body'. Returns a list of (slide_no, action, used, capacity) notes."""
    rules = spec.get("long_content_rules", {})
    demote_ok = rules.get("demote_if_underfull", True)
    prs = open_template(src_path)
    notes = []
    n_long = sum(1 for s in slots if s["is_long"])

    for i, slot in enumerate(slots, start=1):
        if not slot["is_long"]:
            continue
        tmpl_shape = find_shape(prs.slides[slot["tmpl_idx"] - 1], "AI_content")
        usable, cap, max_lines, max_chars, min_lines = long_box_spec(spec, tmpl_shape)
        kept, used, log = generate_long_body(model, slot["section"], usable, cap,
                                             max_lines, max_chars, min_lines)
        slot["long_log"] = log
        underfull = used < min_lines
        # Rule: never leave bullet-sized content in a long box -- demote, unless
        # that would drop the deck below its required minimum.
        if underfull and demote_ok and n_long - 1 >= min_long:
            slot["kind"] = slot.pop("orig_kind", slot["kind"])
            slot["is_long"] = False
            slot["tmpl_idx"] = slot["pre_long_tmpl_idx"]
            slot["section"] = dict(slot["section"], needs_long_content=False)
            n_long -= 1
            notes.append((i, f"demoted to {slot['kind']} (only {used}/{cap} lines)",
                          used, cap))
        else:
            slot["long_body"] = "\n".join(kept)
            tag = "kept (under-full, required)" if underfull else "ok"
            notes.append((i, tag, used, cap))
    return notes


def fill_long_content(slide, section, spec, model, slide_w_in, body=None):
    """Fill a long-content slide (21/22): a short AI_title heading plus AI_content
    as PLAIN multi-line text (no bullets), preserving line breaks and leading
    indentation, trimmed to what the large box can hold. `body` is the block from
    prefetch_long_content(); when None the model is called directly. Returns
    (heading, kept_lines, used_lines, capacity, n_tech)."""
    lc_ph = spec["placeholders"]["long_content"]

    # 1) AI_title FIRST (short heading, <= 25 chars, forced to one line) so we know
    #    exactly how much vertical room it occupies before placing the content.
    heading = truncate(section.get("heading", ""),
                       lc_ph["AI_title"]["max_chars"], ellipsis=True)
    th = find_shape(slide, "AI_title")
    if th is not None:
        fill_simple(th, [heading])
        tf, pt = run_font(th, "Impact", TITLE_PT_DEFAULT)
        fit_title_width(slide, tf, pt, slide_w_in)
        fix_title_geometry(slide, tf, pt)                # one-line height

    # 2) Push the content box clear of the title, THEN measure the (final) box.
    ensure_gap_below_title(slide, LONG_GAP_IN)
    content_shape = find_shape(slide, "AI_content")
    usable, capacity, max_lines, max_chars, min_lines = long_box_spec(
        spec, content_shape)

    # 3) Generate if needed, trim to the box (with an absolute height backstop),
    #    and write as plain multi-line text (fill_simple adds no bullet).
    if body is None:
        body = stage2_long_content(model, section, max_lines, max_chars, min_lines)
        time.sleep(API_CALL_DELAY)
    kept, used, _dropped = measure_long_body(body, usable, capacity)
    kept, used = trim_to_box_height(content_shape, kept)
    fill_simple(content_shape, kept if kept else [""])

    n_tech = sum(1 for ln in kept if _looks_technical(ln))
    return heading, kept, used, capacity, n_tech


# ---------------------------------------------------------------------------
# Orchestration: generate_content (Gemini Stage 1) / fill_placeholders
# (Gemini Stage 2 + write) / save_output
# ---------------------------------------------------------------------------
def generate_content(topic, spec, sequence, model):
    """Gemini Stage 1: one outline call for the whole deck -- title, subtitle,
    wheel tags, and a heading+topic for every Wheel/Bar/Hexagon slide. Returns
    a content dict consumed by fill_placeholders."""
    bar_count = sum(1 for t in sequence if t[0] == "Bars")
    hex_count = sum(1 for t in sequence if t[0] == "Hexagon")
    wheel_count = sum(1 for t in sequence if t[0] == "Wheel")
    content_slide_count = len(sequence) - 2   # minus Start + End

    ph_start = spec["placeholders"]["start"]
    max_title = ph_start["AI_title"]["max_chars"]
    max_subtitle = ph_start["AI_subtitle"]["max_chars"]
    max_tag = spec["placeholders"]["wheel"]["AI_tag_1"]["max_chars"]

    print(f"Topic: {topic!r} | {len(sequence)} slides = "
          f"1 Start + {wheel_count} Wheel + {bar_count} Bar + "
          f"{hex_count} Hexagon + 1 End")

    min_long, max_long = long_content_limits(len(sequence), spec)
    if max_long <= 0:
        guidance = ("This template has no long-content slides: set "
                    "\"needs_long_content\": false on every section.")
    elif min_long > 0:
        guidance = (f"This deck is {len(sequence)} slides -- a LONG deck. Mark AT LEAST "
                    f"{min_long} and at most {max_long} sections "
                    f"\"needs_long_content\": true. Longer decks carry more detailed "
                    f"content, so at least one section must use the large box; aim for "
                    f"{min_long}-{max_long} depending on how technical the topic is.")
    else:
        guidance = (f"This deck is {len(sequence)} slides -- a SHORTER deck. Mark 0-"
                    f"{max_long} sections \"needs_long_content\": true. Zero is "
                    f"perfectly fine for a simple overview topic; do not force a "
                    f"long-content slide if the content works well as bullets.")

    print("\n=== STAGE 1: outline ===")
    print(f"  long-content policy: min={min_long} max={max_long} "
          f"(deck size {len(sequence)})")
    p1 = (STAGE1_PROMPT.replace("{TOPIC}", topic)
          .replace("{content_slide_count}", str(content_slide_count))
          .replace("{bar_count}", str(bar_count))
          .replace("{hex_count}", str(hex_count))
          .replace("{long_content_guidance}", guidance))
    outline = gemini_json(model, p1)
    time.sleep(API_CALL_DELAY)

    # One-line headings: truncate over-limit text with a trailing "..." .
    title = truncate(outline.get("title", ""), max_title, ellipsis=True)
    subtitle = truncate(outline.get("subtitle", ""), max_subtitle, ellipsis=True)
    wheel_tags = [truncate(t, max_tag) for t in
                  (outline.get("wheel_tags") or [])][:wheel_count]
    sections = outline.get("sections") or []
    print(f"  title     : {title!r}  ({len(title)} chars, <= {max_title})")
    print(f"  subtitle  : {subtitle!r}  ({len(subtitle)} chars, <= {max_subtitle})")
    print(f"  wheel_tags: {wheel_tags}")
    print(f"  sections  : {len(sections)} (need {content_slide_count})")
    for s in sections:
        lc = "  [needs_long_content]" if s.get("needs_long_content") else ""
        print(f"     - {s.get('type'):7} heading={s.get('heading')!r} "
              f"topic={s.get('topic')!r}{lc}")

    return {"title": title, "subtitle": subtitle, "wheel_tags": wheel_tags,
            "sections": sections}


def fill_placeholders(prs, content, spec, slots, model, topic=""):
    """Gemini Stage 2 + write everything into `prs`, iterating the routed `slots`
    from route_plan(). Normal Wheel/Bar/Hexagon slots get adaptive, box-fitted
    content (per-point format, fill target per section -- 70-85% by default,
    70-80% for the smaller wheel boxes); LongContent slots get a plain
    multi-line block (code / derivation / paragraph) with line breaks preserved.

    Each slide's content is generated for ONE focus only -- a wheel slide's focus
    is its MIDDLE tag (detected by geometry), a bar/hexagon slide's focus is its
    heading -- and each generation is told what the previous and next slides cover
    so slides stay on-topic and don't overlap.

    Returns a summary list of (slide#, kind, n_items, n_lines, capacity, n_tech)."""
    slide_h_in = prs.slide_height / EMU
    slide_w_in = prs.slide_width / EMU
    marL = content_marL_emu()
    hang = marL / EMU
    wheel_tags = content.get("wheel_tags") or []

    print("\n=== STAGE 2 + FILL (per-slide focus, box-fitted, continuity-aware) ===")

    # --- Per-slot focus/title (drives each slide's content AND the prev/next
    #     continuity context). Wheel focus is the middle tag (geometry-detected);
    #     bar/hexagon/long focus is the heading. Stashed on the slot for reuse. --
    slot_titles = []
    for j, s in enumerate(slots):
        k = s["kind"]
        if k == "Start":
            title = content.get("title", "")
        elif k == "End":
            title = "Summary / conclusion"
        elif k == "Wheel":
            title = wheel_focus_tag(prs.slides[j], wheel_tags)
            s["focus"] = title
        else:                                   # Bars / Hexagon / LongContent
            title = (s.get("section") or {}).get("heading", "")
            s["focus"] = title
        slot_titles.append(title)

    def neighbour(j, step):
        """Focus/title of the slide `step` away from j, skipping Start/End so a
        content slide's continuity refers to real content, not the cover."""
        j += step
        while 0 <= j < len(slots) and slots[j]["kind"] in ("Start", "End"):
            j += step
        return slot_titles[j] if 0 <= j < len(slots) else ""

    # --- Parallel content generation for all normal slides ---------------------
    gen_tasks = {}
    for gi, gslot in enumerate(slots):
        gv = gslot["kind"]
        if gv in ("Start", "End") or gslot.get("is_long"):
            continue
        gslide = prs.slides[gi]
        gkey = _KIND_TO_PH[gv]
        gph = spec["placeholders"][gkey]
        gshape = find_shape(gslide, "AI_content")
        _, _, gcx, gcy = geom_emu(gshape)
        gins = get_ins(gshape)
        gusable = gcx / EMU - gins["l"] - gins["r"]
        gcap = box_line_capacity(gcy, gins)
        gfocus = gslot.get("focus", "")
        if gv == "Wheel":
            # Wheel content is about the MIDDLE tag only -- build a section fixed
            # on that pillar instead of the outline's (possibly all-encompassing)
            # wheel topic, so the box never turns into an overview of every tag.
            gsec = {"type": "wheel", "heading": gfocus, "focus_tag": gfocus,
                    "topic": (f"Introduce and define '{gfocus}', one core concept of "
                              f"{topic}. Cover ONLY {gfocus} -- not the other concepts.")}
        else:
            gsec = gslot["section"]
        gen_tasks[gi] = (gsec, gusable, gph["AI_content"], gcap, gfocus,
                         neighbour(gi, -1), neighbour(gi, +1))

    t_parallel = time.time()
    print(f"  Generating content for {len(gen_tasks)} slides in parallel "
          f"(max_workers=3)...")
    gen_results = {}
    with ThreadPoolExecutor(max_workers=3) as pool:
        futures = {}
        for idx, (sec_p, us_p, phc_p, cap_p, foc_p, prev_p, next_p) in gen_tasks.items():
            f = pool.submit(adaptive_bullets, model, sec_p, us_p, hang,
                            phc_p, cap_p, foc_p, prev_p, next_p)
            futures[f] = idx
        for f in as_completed(futures):
            gen_results[futures[f]] = f.result()
    print(f"  All {len(gen_results)} slides generated in "
          f"{time.time() - t_parallel:.2f}s")

    summary = []
    long_routed = []
    report = []     # per-slide: {slide, type, title, n, unit, fill}
    slide_times = []
    for i, slot in enumerate(slots):
        t_slide = time.time()
        slide = prs.slides[i]
        variant = slot["kind"]

        if variant == "Start":
            _fill_shape_text(slide, "AI_title", [content["title"]])
            _fill_shape_text(slide, "AI_subtitle", [content["subtitle"]])
            print(f"\nSlide {i + 1:2} Start")
            print(f"    title    : {content['title']!r}")
            print(f"    subtitle : {content['subtitle']!r}   (USER_info untouched)")
            report.append({"slide": i + 1, "type": "Start", "title": content["title"],
                           "n": None, "unit": "", "fill": None})
            slide_times.append((i + 1, "Start", time.time() - t_slide))
            continue
        if variant == "End":
            report.append({"slide": i + 1, "type": "End", "title": "-",
                           "n": None, "unit": "", "fill": None})
            slide_times.append((i + 1, "End", time.time() - t_slide))
            continue

        sec = slot["section"]

        # ---- LongContent slot: plain multi-line block, no bullets -----------
        if slot.get("is_long"):
            heading, kept, used, cap, n_tech = fill_long_content(
                slide, sec, spec, model, slide_w_in, body=slot.get("long_body"))
            long_routed.append((i + 1, slot.get("orig_kind", "?"),
                                heading, slot["tmpl_idx"]))
            fill_pct = (100.0 * used / cap) if cap else 0.0
            lc_max = spec["placeholders"]["long_content"]["AI_title"]["max_chars"]
            print(f"\nSlide {i + 1:2} LongContent  "
                  f"(was {slot.get('orig_kind')}, template slide {slot['tmpl_idx']})")
            print(f"    heading={heading!r} ({len(heading)} chars, <= {lc_max})")
            print(f"    -> {len(kept)} text lines -> {used} display lines "
                  f"({fill_pct:.0f}% of {cap}-line box)  tech-lines={n_tech}  "
                  f"[PLAIN TEXT, no bullets, indentation preserved]")
            for ln in kept:
                print(f"        | {ln}")
            summary.append((i + 1, "LongContent", len(kept), used, cap, n_tech))
            report.append({"slide": i + 1, "type": "LongContent", "title": heading,
                           "n": len(kept), "unit": "lines", "fill": fill_pct})
            slide_times.append((i + 1, "LongContent", time.time() - t_slide))
            continue

        # ---- Normal Wheel / Bars / Hexagon slot -----------------------------
        key = _KIND_TO_PH[variant]
        ph = spec["placeholders"][key]

        content_shape = find_shape(slide, "AI_content")
        _, _, cx, cy = geom_emu(content_shape)
        ins = get_ins(content_shape)
        usable = cx / EMU - ins["l"] - ins["r"]
        capacity = box_line_capacity(cy, ins)

        bullets, _tot, log = gen_results[i]

        # headings / wheel tags
        if variant == "Wheel":
            tag_ph = spec["placeholders"]["wheel"]
            focus_n = wheel_focus_index(slide)
            focus_tag = slot.get("focus", "")
            shown = []
            for ti, tag in enumerate(content["wheel_tags"], start=1):
                cfg = tag_ph.get(f"AI_tag_{ti}", tag_ph.get("AI_tag_1", {}))
                tag_lines = wrap_to_lines(tag, cfg.get("max_chars_per_line", 18),
                                          cfg.get("max_lines", 2))
                _fill_shape_text(slide, f"AI_tag_{ti}", tag_lines)
                star = " <focus>" if ti == focus_n else ""
                shown.append(" / ".join(tag_lines) + star)
            head_line = (f"FOCUS (middle tag) = {focus_tag!r} [AI_tag_{focus_n}];  "
                         f"tags={shown}")
        else:
            max_heading = ph["AI_title"]["max_chars"]
            heading = truncate(sec.get("heading", ""), max_heading, ellipsis=True)
            _fill_shape_text(slide, "AI_title", [heading])
            th = find_shape(slide, "AI_title")
            wnote = ""
            if th is not None:
                tf, pt = run_font(th, "Impact", TITLE_PT_DEFAULT)
                winfo = fit_title_width(slide, tf, pt, slide_w_in)   # widen to 1 line
                fix_title_geometry(slide, tf, pt)                    # 1-line height
                if winfo:
                    wnote = (f"  [title box W->{winfo['box_w_in']:.2f}in, "
                             f"+{winfo['widened_in']:.2f}in, "
                             f"one-line={'yes' if winfo['fit'] else 'NO'}]")
            head_line = (f"heading={heading!r} ({len(heading)} chars, "
                         f"<= {max_heading}){wnote}")

        # fill content; height-fit may drop lines -> recount on the FINAL set
        final_bullets, fit_note, _anim = fill_content_bullets(
            slide, bullets, slide_h_in, MIN_TOP_IN, marL)
        final_lines = total_display_lines(final_bullets, usable, hang)
        final_wraps = count_wraps(final_bullets, usable, hang)
        n_tech = sum(1 for b in final_bullets if _looks_technical(b))
        fill_pct = (100.0 * final_lines / capacity) if capacity else 0.0

        print(f"\nSlide {i + 1:2} {variant}")
        print(f"    {head_line}")
        for step in log:
            print(f"    - {step}")
        print(f"    -> FINAL items={len(final_bullets)} lines={final_lines} "
              f"({fill_pct:.0f}% of {capacity}-line box)  wraps={final_wraps}  "
              f"tech-lines={n_tech}  [{fit_note}]  (usable={usable:.2f}in)")
        for b in final_bullets:
            mark = "  <-- formula/code" if _looks_technical(b) else ""
            print(f"        - {b}{mark}")
        summary.append((i + 1, variant, len(final_bullets), final_lines,
                        capacity, n_tech))
        rep_title = (f"{slot.get('focus', '')} (focus)" if variant == "Wheel"
                     else heading)
        report.append({"slide": i + 1, "type": variant, "title": rep_title,
                       "n": len(final_bullets), "unit": "bullets", "fill": fill_pct})
        slide_times.append((i + 1, variant, time.time() - t_slide))

    print("\n=== SUMMARY: fill % and technical-line counts ===")
    for sidx, variant, nb, nlines, cap, ntech in summary:
        pct = (100.0 * nlines / cap) if cap else 0.0
        print(f"   Slide {sidx:2} {variant:11}: items={nb} lines={nlines}/{cap} "
              f"({pct:.0f}% full)  formula/code lines={ntech}")

    print("\n=== PER-SLIDE REPORT ===")
    print(f"   {'#':>2}  {'type':11}  {'title':32}  {'count':>12}  fill")
    print(f"   {'-' * 2}  {'-' * 11}  {'-' * 32}  {'-' * 12}  {'-' * 5}")
    for r in report:
        cnt = "-" if r["n"] is None else f"{r['n']} {r['unit']}"
        pct = "n/a" if r["fill"] is None else f"{r['fill']:.0f}%"
        title = truncate(r["title"] or "", 32)
        print(f"   {r['slide']:>2}  {r['type']:11}  {title:32}  {cnt:>12}  {pct}")

    if long_routed:
        print("\n=== long-content slides routed (21/22) ===")
        for sidx, orig, heading, tmpl in long_routed:
            print(f"   Slide {sidx:2}: {heading!r}  (was {orig}, "
                  f"template slide {tmpl})")
    else:
        print("\n(no long-content slides needed for this deck)")

    print("\n=== PER-SLIDE TIMING ===")
    for sno, stype, dt in sorted(slide_times, key=lambda x: -x[2]):
        print(f"   Slide {sno:>2} {stype:11}  {dt:.2f}s")
    return summary


def save_output(prs, output_path):
    prs.save(output_path)
    print(f"\nSaved: {output_path}")
    return output_path


# Filler words dropped when deriving a short topic slug for the download filename.
# Verbs/articles/meta words a prompt wraps a topic in ("make a presentation about ...",
# "explain the ...") so the slug is the actual subject, not the request boilerplate.
_TOPIC_FILLER = {
    "make", "create", "generate", "build", "design", "prepare", "produce", "write",
    "give", "show", "do", "explain", "describe", "tell", "cover", "covering", "teach",
    "a", "an", "the", "some", "my", "our", "your", "this", "that", "these", "those",
    "presentation", "presentations", "slide", "slides", "deck", "slideshow",
    "powerpoint", "ppt", "pptx", "talk", "lecture", "report", "overview",
    "introduction", "intro", "about", "on", "of", "for", "to", "with", "in", "into",
    "regarding", "concerning", "and", "or", "please", "me", "us", "using", "use",
}


def topic_slug(topic, max_words=2):
    """Short lowercase_underscore slug naming the SUBJECT of a prompt, for filenames.

    Strips request boilerplate ("make a presentation about ...") and keeps the first
    `max_words` significant words, e.g.:
      "make a presentation about machine learning algorithms" -> "machine_learning"
      "The French Revolution causes and consequences"         -> "french_revolution"
      "Digital marketing strategies for small businesses"     -> "digital_marketing"
    Falls back to the first few raw words if everything is filler, and to "deck" if
    the prompt is empty. Never returns special characters."""
    words = re.findall(r"[a-z0-9]+", (topic or "").lower())
    significant = [w for w in words if w not in _TOPIC_FILLER]
    chosen = significant[:max_words] if len(significant) >= 2 else (significant or words)[:3]
    return "_".join(chosen) if chosen else "deck"


def output_filename(topic, template, max_len=40):
    """Download filename '{topic}_{template}.pptx', lowercase with underscores and
    no special characters, capped at `max_len` characters total (extension
    included). The topic slug is trimmed first if the whole name is too long."""
    template = re.sub(r"[^a-z0-9]+", "", (template or "template").lower()) or "template"
    ext = ".pptx"
    slug = topic_slug(topic)
    stem = f"{slug}_{template}"
    if len(stem) + len(ext) > max_len:
        stem = stem[:max_len - len(ext)].rstrip("_") or template
    return f"{stem}{ext}"


def generate_ai_deck(template_name, topic, count, output_path=None):
    """End-to-end: load spec -> Gemini outline (Stage 1) -> route long-content
    sections onto slides 21/22 -> clone the planned slides -> Stage 2 fill
    (adaptive bullets, or plain multi-line blocks for long-content) -> save."""
    total_start = time.time()
    folder = os.path.dirname(os.path.abspath(__file__))
    spec = load_spec(template_name)
    src = spec["_template_path"]
    out = output_path or os.path.join(folder, output_filename(topic, template_name))

    seq = build_slide_sequence(count, spec)     # validates the count against spec

    model = load_gemini(folder)
    t_stage1 = time.time()
    content = generate_content(topic, spec, seq, model)
    print(f"\n[perf] Stage 1 (outline) took {time.time() - t_stage1:.2f}s")

    # Route needs_long_content sections onto the long-content template slides.
    slots = route_plan(seq, content["sections"], spec)
    min_long, max_long = long_content_limits(count, spec)
    n_long = sum(1 for s in slots if s.get("is_long"))
    forced = sum(1 for s in slots if s.get("forced_long"))
    print(f"\n[route] {n_long} long-content slide(s) routed to template "
          f"slides {[s['tmpl_idx'] for s in slots if s.get('is_long')] or '(none)'} "
          f"(policy min={min_long} max={max_long}"
          f"{f', {forced} promoted to meet the minimum' if forced else ''})")

    # Generate the long blocks BEFORE cloning so an under-full one can be demoted
    # back to a normal Bars/Hexagon slide rather than sitting in a huge empty box.
    if n_long:
        t_stage2a = time.time()
        print("\n=== STAGE 2a: long-content blocks (>=60% fill required) ===")
        for sno, action, used, cap in prefetch_long_content(
                slots, spec, model, src, min_long):
            print(f"   slide {sno:2}: {action}  [{used}/{cap} lines]")
        for slot in slots:
            for step in slot.get("long_log", []):
                print(f"      - slot: {step}")
        n_long = sum(1 for s in slots if s.get("is_long"))
        print(f"   -> {n_long} long-content slide(s) after fill check")
        print(f"\n[perf] Stage 2a (long-content prefetch) took {time.time() - t_stage2a:.2f}s")

    t_build = time.time()
    build_deck_from_slots(slots, src, out)
    prs = Presentation(out)
    print(f"\n[perf] Deck cloning + assembly took {time.time() - t_build:.2f}s")

    t_stage2 = time.time()
    fill_placeholders(prs, content, spec, slots, model, topic)
    print(f"\n[perf] Stage 2 (fill all placeholders) took {time.time() - t_stage2:.2f}s")

    t_save = time.time()
    save_output(prs, out)      # save BEFORE reporting: a printing failure must
    print(f"[perf] Save took {time.time() - t_save:.2f}s")
    print_inserted(prs)        # never cost a deck that is already complete

    total_end = time.time()
    print(f"\n{'=' * 62}")
    print(f"[perf] TOTAL EXECUTION TIME: {total_end - total_start:.2f}s")
    print(f"{'=' * 62}")
    return out


# ---------------------------------------------------------------------------
# Offline long-content test (no Gemini calls)
#
# run_build_tests() only walks build_slide_sequence(), so it never clones the
# long-content template slides. This path does: it fabricates an outline with
# two needs_long_content Bar sections, routes them through the real route_plan()
# / build_deck_from_slots() / fill_long_content() code, and feeds Stage 2 from a
# canned stub instead of the API.
# ---------------------------------------------------------------------------
_STUB_CODE_BLOCK = (
    "def fizzbuzz(n):\n"
    "    results = []\n"
    "    for i in range(1, n + 1):\n"
    "        if i % 15 == 0:\n"
    "            results.append(\"FizzBuzz\")\n"
    "        elif i % 3 == 0:\n"
    "            results.append(\"Fizz\")\n"
    "        elif i % 5 == 0:\n"
    "            results.append(\"Buzz\")\n"
    "        else:\n"
    "            results.append(str(i))\n"
    "    return results\n"
    "\n"
    "print(fizzbuzz(20))"
)

_STUB_PARAGRAPH = (
    "A loop is a control structure that repeats a block of statements until a\n"
    "condition stops holding. The for-loop walks a known sequence and is the right\n"
    "choice when the number of iterations is fixed in advance; the while-loop keeps\n"
    "going while a condition stays true and suits work whose length is not known\n"
    "until it runs. Both can be exited early with break or advanced with continue.\n"
    "\n"
    "Choosing between them is a readability decision, not a performance one. If the\n"
    "reader can see the bounds at a glance, prefer the for-loop; reach for while only\n"
    "when the stopping condition is genuinely dynamic."
)


class _StubResponse:
    def __init__(self, text):
        self.text = text


class _StubModel:
    """Stands in for _FallbackModel in offline tests. Returns canned Stage 2
    long-content JSON, picked by which heading appears in the prompt, and records
    every call so the test can assert that zero real API traffic was needed."""

    name = "stub (offline, no API calls)"

    def __init__(self, by_heading):
        self._by_heading = by_heading
        self.calls = []

    def generate_content(self, prompt):
        for heading, body in self._by_heading.items():
            if heading in prompt:
                self.calls.append(heading)
                return _StubResponse(json.dumps({"content": body}))
        self.calls.append("(unmatched)")
        return _StubResponse(json.dumps({"content": "stub content"}))


def run_long_content_test(template_name="neon", output_path=None, count=8):
    """Build a deck whose Bar sections are flagged needs_long_content so slides
    21 and 22 are exercised end-to-end with dummy content -- a code block on 21
    and a prose paragraph on 22 -- without contacting Gemini. Non-long slides get
    static placeholder bullets through the normal fill path. Returns the output
    path."""
    global API_CALL_DELAY
    folder = os.path.dirname(os.path.abspath(__file__))
    spec = load_spec(template_name)
    src = spec["_template_path"]
    out = output_path or os.path.join(folder, "output_longcontent_test.pptx")

    if "long_content" not in spec["sections"]:
        raise ValueError(f"template {template_name!r} has no long_content section")

    headings = ["FizzBuzz in Python", "What a loop is"]
    model = _StubModel({headings[0]: _STUB_CODE_BLOCK,
                        headings[1]: _STUB_PARAGRAPH})

    # --- fabricate the outline Stage 1 would have returned -------------------
    seq = build_slide_sequence(count, spec)
    sections, n_flagged = [], 0
    for kind, _idx, _label in seq:
        if kind in ("Start", "End"):
            continue
        stype = {"Wheel": "wheel", "Bars": "bar", "Hexagon": "hexagon"}[kind]
        flag = stype == "bar" and n_flagged < len(headings)
        sections.append({
            "type": stype,
            "heading": headings[n_flagged] if flag else f"Test {stype} slide",
            "topic": ("show the code" if flag and n_flagged == 0 else
                      "explain loops in prose" if flag else
                      f"placeholder {stype} content"),
            "needs_long_content": flag,
        })
        if flag:
            n_flagged += 1
    content = {"title": "Loops", "subtitle": "Offline long-content test",
               "wheel_tags": ["Definition", "For loop", "While loop"],
               "sections": sections}

    print("=" * 64)
    print(f"LONG-CONTENT TEST ({count} slides, template {template_name!r}, no API)")
    print("=" * 64)

    slots = route_plan(seq, sections, spec)
    long_idx = [s["tmpl_idx"] for s in slots if s["is_long"]]
    print(f"[route] flagged sections: {n_flagged}; "
          f"routed onto template slides {long_idx or '(none)'}")
    for i, slot in enumerate(slots, start=1):
        mark = "  <-- LONG" if slot["is_long"] else ""
        print(f"   slide {i:2}  {slot['label']:3} -> {slot['kind']:11} "
              f"(template {slot['tmpl_idx']}){mark}")
    if not long_idx:
        raise AssertionError("no sections were routed to long-content slides")

    build_deck_from_slots(slots, src, out)
    prs = Presentation(out)

    # --- fill: long slots via the real path, others with static bullets ------
    slide_h_in = prs.slide_height / EMU
    slide_w_in = prs.slide_width / EMU
    marL = content_marL_emu()
    delay, API_CALL_DELAY = API_CALL_DELAY, 0.0     # no API -> no rate limiting
    try:
        for i, slot in enumerate(slots):
            slide = prs.slides[i]
            if slot["kind"] == "Start":
                _fill_shape_text(slide, "AI_title", [content["title"]])
                _fill_shape_text(slide, "AI_subtitle", [content["subtitle"]])
                continue
            if slot["kind"] == "End":
                continue

            sec = slot["section"]
            if slot["is_long"]:
                heading, kept, used, cap, n_tech = fill_long_content(
                    slide, sec, spec, model, slide_w_in)
                pct = (100.0 * used / cap) if cap else 0.0
                print(f"\nSlide {i + 1:2} LongContent "
                      f"(was {slot.get('orig_kind')}, template {slot['tmpl_idx']})")
                print(f"    heading={heading!r}")
                print(f"    -> {len(kept)} text lines -> {used} display lines "
                      f"({pct:.0f}% of {cap}-line box)  tech-lines={n_tech}")
                for ln in kept:
                    print(f"        | {ln}")
                continue

            # plain placeholder bullets, written through the normal fill path
            if slot["kind"] == "Wheel":
                tag_ph = spec["placeholders"]["wheel"]
                for ti, tag in enumerate(content["wheel_tags"], start=1):
                    cfg = tag_ph.get(f"AI_tag_{ti}", tag_ph.get("AI_tag_1", {}))
                    _fill_shape_text(slide, f"AI_tag_{ti}",
                                     wrap_to_lines(tag, cfg.get("max_chars_per_line", 18),
                                                   cfg.get("max_lines", 2)))
            else:
                _fill_shape_text(slide, "AI_title", [sec["heading"]])
                th = find_shape(slide, "AI_title")
                if th is not None:
                    tf, pt = run_font(th, "Impact", TITLE_PT_DEFAULT)
                    fit_title_width(slide, tf, pt, slide_w_in)
                    fix_title_geometry(slide, tf, pt)
            bullets = [f"Placeholder bullet {b}" for b in range(1, 4)]
            fill_content_bullets(slide, bullets, slide_h_in, MIN_TOP_IN, marL)
            print(f"\nSlide {i + 1:2} {slot['kind']:11} (placeholder bullets)")
    finally:
        API_CALL_DELAY = delay

    print(f"\n[stub] Gemini calls made: 0 real, "
          f"{len(model.calls)} served from the stub {model.calls}")
    save_output(prs, out)
    return out


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------
def parse_args():
    p = argparse.ArgumentParser(
        description="Fill a spec-driven PowerPoint template with AI-generated content.")
    p.add_argument("--template", default="neon",
                    help="Template name (folder under templates/, e.g. 'neon')")
    p.add_argument("--topic", help="Presentation topic (required unless --test-long-content)")
    p.add_argument("--count", type=int, help="Target slide count")
    p.add_argument("--output", default=None,
                    help="Output .pptx path (default: output_<template>.pptx)")
    p.add_argument("--test-long-content", action="store_true",
                    help="Offline test: build a deck exercising the long-content "
                         "slides with dummy content (no Gemini calls)")
    args = p.parse_args()
    if not args.test_long_content and (not args.topic or not args.count):
        p.error("--topic and --count are required unless --test-long-content is set")
    return args


if __name__ == "__main__":
    args = parse_args()
    try:
        if args.test_long_content:
            run_long_content_test(args.template, args.output,
                                  count=args.count or 8)
        else:
            generate_ai_deck(args.template, args.topic, args.count, args.output)
    except Exception as exc:  # noqa: BLE001 - surface the failure clearly
        import traceback
        print(f"\n[ERROR] {type(exc).__name__}: {exc}")
        traceback.print_exc()
